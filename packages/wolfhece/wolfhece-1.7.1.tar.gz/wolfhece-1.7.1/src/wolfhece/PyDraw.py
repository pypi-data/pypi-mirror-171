import numpy as np
from wx import dataview, TreeCtrl
import wx
import wx.propgrid as pg
from wx.core import VERTICAL, BoxSizer, Height, ListCtrl, StaticText, TextCtrl, Width
from wx.glcanvas import GLCanvas, GLContext
from wx.dataview import TreeListCtrl
import wx.lib.ogl as ogl
from PIL import Image
from PIL.PngImagePlugin import PngInfo
from io import BytesIO
from OpenGL.GL import *
from OpenGL.GLUT import *
import matplotlib.pyplot as plt
from matplotlib.widgets import Button as mplButton
from os import scandir, listdir
from os.path import exists, join, normpath
from pptx import Presentation
import threading

from .wolf_texture import genericImagetexture,imagetexture,Text_Image_Texture
from .xyz_file import xyz_scandir, XYZFile
from .mesh2d import wolf2dprev
from .PyPalette import wolfpalette
from .wolfresults_2D import Wolfresults_2D, CHOICES_VIEW_2D
from .PyTranslate import _
from .PyVertex import cloud_vertices, getIfromRGB
from .RatingCurve import SPWMIGaugingStations, SPWDCENNGaugingStations
from .wolf_array import WOLF_ARRAY_MB, SelectionData, WolfArray, WolfArray_Sim2D, WolfArrayMB, CropDialog, header_wolf
from .PyParams import Wolf_Param
from .BcManager import BcManager
from .PyVertexvectors import *
from .Results2DGPU import wolfres2DGPU
from .PyCrosssections import crosssections, profile
from .GraphNotebook import PlotNotebook
from .laz_viewer import myviewer, read_laz, clip_data_xyz
from . import Lidar2002

ID_SELECTCS = 1000
ID_SORTALONG = 1001
ID_LOCMINMAX = 1002



class WolfMapViewer(wx.Frame):
    """
    Fenêtre de visualisation de données WOLF grâce aux WxWidgets
    """

    TIMER_ID = 100  # délai d'attente avant action

    mybc: BcManager = None  # Gestionnaire de CL
    myarrays: list  # matrices ajoutées
    myvectors: list  # zones vectorielles ajoutées
    myclouds: list  # nuages de vertices
    myothers: list
    mywmsback: list
    mywmsfore: list
    myres2D: list

    canvas: ogl.ShapeCanvas # GLCanvas  # canvas OpenGL
    context: GLContext  # context OpenGL
    mytooltip: Wolf_Param  # Objet WOLF permettant l'analyse de ce qui est sous la souris
    treelist: TreeListCtrl  # Gestion des éléments sous forme d'arbre
    selitem: StaticText
    leftbox: BoxSizer
    added: dict  # dictionnaire des éléments ajoutés

    active_vector: vector
    active_zone: zone
    active_zones: Zones
    active_array: WolfArray
    active_vertex: wolfvertex
    active_cs: crosssections

    def __init__(self, wxparent, title, w=500, h=500, treewidth=200, wolfparent=None):

        self.action = None  # Action à entreprendre
        self.update_absolute_minmax = False  # Force la MAJ de la palette
        self.copyfrom = None  # aucun élément pointé par CTRL+C
        
        self.wolfparent = wolfparent

        self.regular = True  # Gestion de la taille de fenêtre d'affichage, y compris l'arbre de gestion
        self.sx = 1  # facteur d'échelle selon X = largeur en pixels/largeur réelle
        self.sy = 1  # facteur d'échelle selon Y = hauteur en pixels/hauteur réelle
        self.samescale = True  # force le même facteur d'échelle

        self.dynapar_dist = 1.

        # emprise initiale
        self.xmin = 0.
        self.ymin = 0.
        self.xmax = 40.
        self.ymax = 40.
        self.width = self.xmax - self.xmin  # largeur de la zone d'affichage en coordonnées réelles
        self.height = self.ymax - self.ymin  # hauteur de la zone d'affichage en coordonnées réelles
        self.canvaswidth = 100
        self.canvasheight = 100

        # position de la caméra
        self.mousex = self.width / 2.
        self.mousey = self.height / 2.

        self.bordersize = 0  # zone réservée au contour
        self.titlesize = 0  # zone réservée au titre
        self.treewidth = 200  # largeur de la zone d'arbre "treelist"

        self.backcolor = wx.Colour(255, 255, 255)  # couleur de fond
        self.mousedown = (0., 0.)  # position initiale du bouton position bas
        self.mouseup = (0., 0.)  # position initiale du bouton position haut
        self.oneclick = True  # détection d'un simple click ou d'un double-click
        self.move = False  # la souris est-elle en train de bouger?

        self.linked = False
        self.link_shareopsvect = True
        self.linkedList = None
        self.link_params = None

        self.forcemimic = True
        self.currently_readresults = False

        self.mylazdata = None

        self.treewidth = treewidth
        super(WolfMapViewer, self).__init__(wxparent, title=title, size=(w + self.treewidth, h))

        # Gestion des menus
        self.popupmenu = wx.Menu()
        self.popupmenu.Bind(wx.EVT_MENU, self.OnPopupItemSelected)
        for text in ['Save', 'Save as']:
            item = self.popupmenu.Append(-1, text)

        self.menubar = wx.MenuBar()

        self.menuwolf2d = None
        self.menusim2D = None

        self.filemenu = wx.Menu()
        openitem = self.filemenu.Append(wx.ID_OPEN, _('Open project'), _('Open project'))
        saveproject = self.filemenu.Append(wx.ID_ANY, _('Save project'), _('save project'))
        self.filemenu.AppendSeparator()
        saveitem = self.filemenu.Append(wx.ID_SAVE, _('Save'), _('Save data to files'))
        saveasitem = self.filemenu.Append(wx.ID_SAVEAS, _('Save as...'), _('Save data to new files'))
        savecanvas = self.filemenu.Append(wx.ID_ANY, _('Save to image...'), _('Save canvas to image file'))

        self.filemenu.AppendSeparator()
        # --- GLTF
        self.menugltf = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Gltf2...'), self.menugltf)
        
        exportgltf = self.menugltf.Append(wx.ID_ANY, _('Export to gltf...'), _('Save data to gltf files'))
        importgltf = self.menugltf.Append(wx.ID_ANY, _('Import from gltf...'), _('Import data from gltf files'))
        compareitem = self.menugltf.Append(wx.ID_ANY, _('Set gltf comparison'), _('Set gltf comparison'))
        updategltf = self.menugltf.Append(wx.ID_ANY, _('Update from gltf...'), _('Update data from gltf files'))
        
        self.filemenu.AppendSeparator()
        compareitem = self.filemenu.Append(wx.ID_ANY, _('Set comparison'), _('Set comparison'))
        multiview = self.filemenu.Append(wx.ID_ANY, _('Multiviewer'), _('Multiviewer'))
        self.filemenu.AppendSeparator()


        # ---
        self.menucreateobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Create...'),self.menucreateobj)

        createarray = self.menucreateobj.Append(wx.ID_FILE6, _('Create array...'), _('New array (binary file - real)'))
        createarray2002 = self.menucreateobj.Append(wx.ID_ANY, _('Create array from Lidar 2002...'),
                                               _('Create array from Lidar 2002 (binary file - real)'))
        createarrayxyz = self.menucreateobj.Append(wx.ID_ANY, _('Create array from bathymetry file...'),
                                              _('Create array from XYZ (ascii file - real)'))
        createvector = self.menucreateobj.Append(wx.ID_FILE7, _('Create vectors...'), _('New vectors'))
        createcloud = self.menucreateobj.Append(wx.ID_FILE8, _('Create cloud...'), _('New cloud'))
        self.filemenu.AppendSeparator()


        # -----        
        self.menuaddobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Add...'),self.menuaddobj)
        
        addarray = self.menuaddobj.Append(wx.ID_FILE1, _('Add array...'), _('Add array (binary file - real)'))
        addarraycrop = self.menuaddobj.Append(wx.ID_ANY, _('Add array and crop...'),
                                            _('Add array and crop (binary file - real)'))
        addvector = self.menuaddobj.Append(wx.ID_FILE2, _('Add vectors...'), _('Add vectors'))
        addcloud = self.menuaddobj.Append(wx.ID_FILE3, _('Add cloud...'), _('Add cloud'))
        addprofiles = self.menuaddobj.Append(wx.ID_FILE4, _('Add cross sections...'), _('Add cross sections'))
        addres2D = self.menuaddobj.Append(wx.ID_ANY, _('Add Wolf2D results...'), _('Add Wolf 2D results'))
        
        self.filemenu.AppendSeparator()
        addscan = self.filemenu.Append(wx.ID_FILE5, _('Recursive scan...'), _('Add recursively'))

        self.toolsmenu = wx.Menu()
        self.link_cs_zones = self.toolsmenu.Append(wx.ID_ANY, _("Link cross sections to active zones"),
                                                   _("Link cross section"))
        self.sortalong = self.toolsmenu.Append(ID_SORTALONG, _("Sort along..."),
                                               _("Sort cross sections along support vector"))
        self.select_cs = self.toolsmenu.Append(ID_SELECTCS, _("Pick one cross section"), _("Select cross section"),
                                               kind=wx.ITEM_CHECK)
        self.menumanagebanks = self.toolsmenu.Append(wx.ID_ANY, _("Manage banks..."), _("Manage banks"))
        self.menucreatenewbanks = self.toolsmenu.Append(wx.ID_ANY, _("Create banks from vertices..."),
                                                        _("Manage banks"))
        self.renamecs = self.toolsmenu.Append(wx.ID_ANY, _("Rename cross sections..."), _("Rename"))
        self.menutrianglecs = self.toolsmenu.Append(wx.ID_ANY, _("Triangulate cross sections..."), _("Triangulate"))
        self.menuexportgltfonebyone = self.toolsmenu.Append(wx.ID_ANY, _("Export cross sections to gltf..."),
                                                            _("Export gltf"))
        self.menupontgltfonebyone = self.toolsmenu.Append(wx.ID_ANY, _("Create bridge and export gltf..."),
                                                          _("Bridge gltf"))
        self.menuviewerinterpcs = None
        self.menuinterpcs = None

        self.minmaxmenu = wx.Menu()
        self.locminmax = self.minmaxmenu.Append(ID_LOCMINMAX, _("Local minmax"), _("Adapt palette on current zoom"),
                                                kind=wx.ITEM_CHECK)

        self.lazmenu = wx.Menu()
        croplaz = self.lazmenu.Append(wx.ID_ANY, _('Clip LAZ data on current zoom'), _('Clip LAZ'))
        readlaz = self.lazmenu.Append(wx.ID_ANY, _('Read LAZ data from npz'), _('read LAZ'))
        viewlaz = self.lazmenu.Append(wx.ID_ANY, _('Create LAZ viewer'), _('LAZ Viewer'))
        bridgelaz = self.lazmenu.Append(wx.ID_ANY, _('Create cloud points from bridges'), _('LAZ Bridge'))
        buildinglaz = self.lazmenu.Append(wx.ID_ANY, _('Create cloud points from buildings'), _('LAZ Buildings'))

        self.analyzemenu = wx.Menu()
        plotvect = self.analyzemenu.Append(wx.ID_ANY, _("Plot active vector..."),
                                           _("Plot the active vector and linked arrays"))
        plotpoly = self.analyzemenu.Append(wx.ID_ANY, _("Plot active polygons..."),
                                           _("Plot the active polygons and linked arrays"))

        # Gestion des outils --> Utile pour ManageActions
        self.tools = {}
        curtool = self.tools[ID_SELECTCS] = {}
        curtool['menu'] = self.select_cs
        curtool['name'] = 'Select nearest profile'

        self.active_vector = None
        self.active_zone = None
        self.active_zones = None
        self.active_vertex = None
        self.active_array = None
        self.active_cs = None
        self.active_res2d = None

        curtool = self.tools[ID_SORTALONG] = {}
        curtool['menu'] = self.sortalong
        curtool['name'] = 'Sort along vector'

        curtool = self.tools[ID_LOCMINMAX] = {}
        curtool['menu'] = self.locminmax
        curtool['name'] = None

        self.menubar.Append(self.filemenu, _('&File'))
        self.menubar.Append(self.toolsmenu, _('&Cross sections'))
        self.menubar.Append(self.lazmenu, _('&LAZ data'))
        self.menubar.Append(self.minmaxmenu, _('&Palette'))
        self.menubar.Append(self.analyzemenu, _('&Analyze'))
        self.SetMenuBar(self.menubar)
        self.Bind(wx.EVT_MENU, self.OnMenubar)

        # Ajout du conteneur OpenGL
        self.canvas = GLCanvas(self)
        self.context = GLContext(self.canvas)
        self.mybackisloaded = False
        self.myfrontisloaded = False

        # ajout d'une liste en arbre des objets
        self.treelist = TreeListCtrl(self, style=wx.dataview.TL_CHECKBOX | wx.LC_EDIT_LABELS)
        self.selitem = StaticText(self, style=wx.ALIGN_CENTER_HORIZONTAL)
        self.selobj = None

        self.root = self.treelist.GetRootItem()
        self.treelist.AppendColumn(_('Objects to plot'))
        self.myitemsarray = self.treelist.AppendItem(self.root, _("Arrays"))
        self.myitemsvector = self.treelist.AppendItem(self.root, _("Vectors"))
        self.myitemscloud = self.treelist.AppendItem(self.root, _("Clouds"))
        self.myitemsres2d = self.treelist.AppendItem(self.root, _("Wolf2D"))
        self.myitemsothers = self.treelist.AppendItem(self.root, _("Others"))
        self.myitemswmsback = self.treelist.AppendItem(self.root, _("WMS-background"))
        self.myitemswmsfore = self.treelist.AppendItem(self.root, _("WMS-foreground"))

        self.added = {}
        self.added['arrays'] = {}
        self.added['vectors'] = {}
        self.added['clouds'] = {}
        self.added['wolf2d'] = {}
        self.added['others'] = {}
        self.added['wms-background'] = {}
        self.added['wms-foreground'] = {}

        width, height = self.GetClientSize()
        self.bordersize = int((w - width + self.treewidth) / 2)
        self.titlesize = h - height - self.bordersize
        self.SetSize(w + 2 * self.bordersize + self.treewidth, h + self.bordersize + self.titlesize)

        # dimensionnement et positionnement de la fenêtre OpenGL
        self.canvas.SetSize(width - self.treewidth, height)
        self.canvas.SetPosition((self.treewidth, 0))

        self.setbounds()

        # dimensionnement et positionnement de l'arbre
        self.leftbox = BoxSizer(orient=wx.VERTICAL)
        self.leftbox.Add(self.treelist, 1, wx.LEFT)
        self.leftbox.Add(self.selitem, 0, wx.LEFT)
        self.treelist.SetSize(self.treewidth, height)
        # self.selitem.SetSize(self.treewidth,30)
        self.SetSizer(self.leftbox)

        # self.treelist.SetPosition((0,0))
        # self.selitem.SetPosition((0,height-30))

        # fenêtre ToolTip
        self.mytooltip = Wolf_Param(self, "Values", ontop=True, to_read=False, withbuttons=False)
        self.mytooltip.SetSize(300, 400)
        self.mytooltip.Show(False)

        self.notebookcs = None
        self.notebookbanks = None
        self.myaxcs = None
        self.myfigcs = None

        self.InitUI()

    def get_choices_arrays(self):
        
        mychoices = wx.Choice(self,wx.ID_ANY,choices=[cur.idx for cur in self.myarrays])


    def menu_wolf2d(self):

        if self.menuwolf2d is None:
            self.menuwolf2d = wx.Menu()
            self.menu2d_curentview = self.menuwolf2d.Append(wx.ID_ANY, _("Change current view"), _("Current view"))
            self.menu2d_lastres = self.menuwolf2d.Append(wx.ID_ANY, _("Read last result"), _("Current view"))
            self.menu2d_bc = self.menuwolf2d.Append(wx.ID_ANY, _("Manage boundary conditions..."), _("BC manager"))
            # self.menu2d_tft_bc = self.menuwolf2d.Append(wx.ID_ANY,_("Transfer boundary conditions..."),_("Transfer BC"))

            self.menubar.Append(self.menuwolf2d, _('Options 2D'))
            
    def menu_sim2D(self):
        if self.menusim2D is None:
            self.menusim2D = wx.Menu()
            self.menubar.Append(self.menusim2D, _('Tools 2D'))        

            update = self.menusim2D.Append(wx.ID_ANY, _('Update model from current mask'), _('Update model'))
            updateblocfile = self.menusim2D.Append(wx.ID_ANY, _('Update .bloc file'), _('Update bloc'))
            updatefreesurface = self.menusim2D.Append(wx.ID_ANY, _('Update free surface elevation - IC'), _('Update free surface elevation'))
            updaterough = self.menusim2D.Append(wx.ID_ANY, _('Update roughness coeff'), _('Update roughness coefficient'))

    def triangulate_cs(self):

        msg = ''
        if self.active_zones is None:
            msg += _(' The active zones is None. Please activate the desired object !\n')
        if self.active_cs is None:
            msg += _(' The is no cross section. Please active the desired object or load file!')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        dlg = wx.NumberEntryDialog(None, _('What is the desired size [cm] ?'), 'ds', 'ds size', 100, 1, 10000.)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        ds = float(dlg.GetValue()) / 100.
        dlg.Destroy()

        self.myinterp = Interpolators(self.active_zones, self.active_cs, ds)

        self.add_object('vector', newobj=self.myinterp.myzones, ToCheck=False, id='Interp_mesh')

        if self.menuviewerinterpcs is None:
            self.menuviewerinterpcs = self.toolsmenu.Append(wx.ID_ANY, _("New cloud Viewer..."),
                                                            _("Cloud viewer Interpolate"))
        if self.menuinterpcs is None:
            self.menuinterpcs = self.toolsmenu.Append(wx.ID_ANY, _("Interpolate on active array..."), _("Interpolate"))

    def interpolate_cs(self):
        if self.active_array is not None and self.myinterp is not None:

            choices = ["nearest", "linear", "cubic"]
            dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            self.myinterp.interp_on_array(self.active_array, method)

    def save_canvasogl(self, fn='', mpl=True, ds=0.):

        self.Paint()

        if fn == '':
            dlg = wx.FileDialog(None, _('Choose file name'), wildcard='PNG (*.png)|*.png|JPG (*.jpg)|*.jpg',
                                style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            fn = dlg.GetPath()
            dlg.Destroy()

        if self.canvas.SetCurrent(self.context):
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage = Image.frombuffer("RGBA", (self.canvaswidth, self.canvasheight), data)
            myimage = myimage.transpose(1)

            metadata = PngInfo()
            metadata.add_text('xmin', str(self.xmin))
            metadata.add_text('ymin', str(self.ymin))
            metadata.add_text('xmax', str(self.xmax))
            metadata.add_text('ymax', str(self.ymax))

            if mpl:
                if ds == 0.:
                    dlg = wx.NumberEntryDialog(self,
                                               _("xmin : {:.3f} \nxmax : {:.3f} \nymin : {:.3f} \nymax : {:.3f} \n\n  dx : {:.3f}\n  dy : {:.3f}").format(
                                                   self.xmin, self.xmax, self.ymin, self.ymax, self.xmax - self.xmin,
                                                   self.ymax - self.ymin),
                                               _("Interval [m]"), _("Ticks interval ?"), 500, 1, 10000)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    ds = float(dlg.GetValue())
                    dlg.Destroy()

                extent = (self.xmin, self.xmax, self.ymin, self.ymax)
                fig, ax = plt.subplots(1, 1)
                ax.imshow(myimage, origin='upper',
                          extent=extent)

                x1 = np.ceil((self.xmin // ds) * ds)
                if x1 < self.xmin:
                    x1 += ds
                x2 = int((self.xmax // ds) * ds)
                if x2 > self.xmax:
                    x2 -= ds
                y1 = np.ceil((self.ymin // ds) * ds)
                if y1 < self.ymin:
                    y1 += ds
                y2 = int((self.ymax // ds) * ds)
                if y2 > self.ymax:
                    y2 -= ds

                x_label_list = np.linspace(x1, x2, int((x2 - x1) / ds) + 1, True)
                x_label_list = np.insert(x_label_list, 0, self.xmin)
                x_label_list = np.insert(x_label_list, -1, self.xmax)
                x_label_list = np.unique(x_label_list)

                y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                y_label_list = np.insert(y_label_list, 0, self.ymin)
                y_label_list = np.insert(y_label_list, -1, self.ymax)
                y_label_list = np.unique(y_label_list)

                ax.set_xticks(x_label_list)
                ax.set_yticks(y_label_list)

                ax.set_xticklabels(plt.FormatStrFormatter('%.1f').format_ticks(x_label_list), fontsize=8, rotation=30)
                ax.set_yticklabels(plt.FormatStrFormatter('%.1f').format_ticks(y_label_list), fontsize=8)

                ax.set_xlabel('X [m]')
                ax.set_ylabel('Y [m]')

                plt.savefig(fn, dpi=150)
            else:
                myimage.save(fn, pnginfo=metadata)

            return fn, ds
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def reporting(self, dir=''):
        if dir == '':
            dlg = wx.DirDialog(None, "Choose directory to store reporting", style=wx.FD_SAVE)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            dir = dlg.GetPath()
            dlg.Destroy()

        myppt = Presentation(__file__)
        slide = myppt.slides.add_slide(0)

        for curzone in self.myzones:
            for curvec in curzone.myvectors:
                curvec: vector
                if curvec.nbvertices > 1:
                    oldwidth = curvec.myprop.width
                    curvec.myprop.width = 4
                    myname = curvec.myname

                    self.Activate_vector(curvec)

                    if self.parentGUI.linked:
                        for curview in self.parentGUI.linkedList:
                            title = curview.GetTitle()
                            curview.zoomon_activevector()
                            fn = path.join(dir, title + '_' + myname + '.png')
                            curview.save_canvasogl(fn)
                    else:
                        self.parentGUI.zoomon_activevector()
                        fn = path.join(dir, myname + '.png')
                        self.parentGUI.save_canvasogl(fn)

                        fn = path.join(dir, 'palette_v_' + myname + '.png')
                        self.parentGUI.active_array.mypal.export_image(fn, 'v')
                        fn = path.join(dir, 'palette_h_' + myname + '.png')
                        self.parentGUI.active_array.mypal.export_image(fn, 'h')

                    curvec.myprop.width = oldwidth

    def InitUI(self):

        self.Bind(wx.EVT_SIZE, self.OnSize)
        self.Bind(wx.EVT_CLOSE, self.OnClose)

        # self.canvas.Bind(wx.EVT_CONTEXT_MENU, self.OnShowPopup)
        self.canvas.Bind(wx.EVT_PAINT, self.OnPaint)
        
        self.treelist.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)
        self.canvas.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)
        
        # self.treelist.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)
        # self.canvas.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)

        self.canvas.Bind(wx.EVT_BUTTON, self.OnButton)
        self.canvas.Bind(wx.EVT_RIGHT_DCLICK, self.OnRDClick)
        self.canvas.Bind(wx.EVT_LEFT_DCLICK, self.OnLDClick)
        self.canvas.Bind(wx.EVT_LEFT_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_MIDDLE_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_RIGHT_DOWN, self.OnRightDown)
        self.canvas.Bind(wx.EVT_RIGHT_UP, self.OnRightUp)
        self.canvas.Bind(wx.EVT_MOTION, self.OnMotion)
        self.canvas.Bind(wx.EVT_LEAVE_WINDOW, self.OnLeave)
        self.canvas.Bind(wx.EVT_MOUSEWHEEL, self.OnButton)

        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CHECKED, self.OnCheckItem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_ACTIVATED, self.OnActivateTreeElem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CONTEXT_MENU, self.OntreeRight)
        # dispo dans wxpython 4.1 self.Bind(wx.EVT_GESTURE_ZOOM,self.OnZoomGesture)

        self.Centre()

        self.myarrays = []
        self.myvectors = []
        self.myclouds = []
        self.myothers = []
        self.mywmsback = []
        self.mywmsfore = []
        self.myres2D = []

        self.Show(True)

    def OnSize(self, e):
        """
        Redimensionnement de la fenêtre
        """
        if self.regular:
            # retrouve la taille de la fenêtre
            width, height = self.GetClientSize()
            # enlève la barre d'arbre
            width -= self.treewidth
            # définit la taille de la fenêtre graphique OpenGL et sa position (à droite de l'arbre)
            self.canvas.SetSize(width, height)
            self.canvas.SetPosition((self.treewidth, 0))
            # calcule les limites visibles sur base de la taille de la fenêtre et des coefficients sx sy
            self.setbounds()
            # fixe la taille de l'arbre (notamment la hauteur)
            # self.treelist.SetSize(self.treewidth,height)
            e.Skip()

    def ManageActions(self, id):

        curmenu = self.tools[id]['menu']

        if curmenu.IsCheckable():
            if not curmenu.IsChecked():
                curmenu.Check(False)
                self.action = None

                if id == ID_LOCMINMAX:
                    self.update_absolute_minmax = True
            else:
                curmenu.Check()
                if not self.tools[id]['name'] is None:
                    self.action = self.tools[id]['name']

        else:
            if id == ID_SORTALONG:
                # Tri le long d'un vecteur
                if not self.active_cs is None and not self.active_vector is None:
                    self.active_cs: crosssections
                    self.active_vector: vector
                    self.active_cs.sort_along(self.active_vector.asshapely_ls(), self.active_vector.myname, False)
                else:
                    msg = ''
                    if self.active_cs is None:
                        msg += _('Please select the active cross sections \n')
                    if self.active_vector is None:
                        msg += _('Please select the active supprt vector')
                    mydiag = wx.MessageDialog(self, msg, _('Sort along'))
                    mydiag.ShowModal()

    def setbounds(self,updatescale=True):
        '''
        Calcule les limites visibles de la fenêtrte graphique sur base des facteurs d'échelle courants
        '''
        
        if updatescale:
            self.updatescalefactors()
        
            # retrouve la taille de la fenêtre OpenGL
            width, height = self.canvas.GetSize()
            self.canvaswidth = width
            self.canvasheight = height

            # calcule la taille selon X et Y en coordonnées réelles
            width = width / self.sx
            height = height / self.sy

            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - width / 2.
            self.xmax = self.xmin + width
            self.ymin = self.mousey - height / 2.
            self.ymax = self.ymin + height

            self.width = width
            self.height = height

            self.mousex = self.xmin + width / 2.
            self.mousey = self.ymin + height / 2.

            self.updatescalefactors()

        else:
            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - self.width / 2.
            self.xmax = self.xmin + self.width
            self.ymin = self.mousey - self.height / 2.
            self.ymax = self.ymin + self.height

        self.mybackisloaded = False
        self.myfrontisloaded = False

        self.Refresh()
        self.mimicme()

    def updatescalefactors(self):
        width, height = self.canvas.GetSize()

        self.sx = 1
        self.sy = 1
        if width > 0:
            self.sx = float(width) / self.width
        if height > 0:
            self.sy = float(height) / self.height

        self.sx = min(self.sx, self.sy)
        self.sy = self.sx

    def add_viewer_and_link(self):
        dlg = wx.TextEntryDialog(self, _('Enter a caption for the new window'))

        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        newcap = dlg.GetValue()
        dlg.Destroy()
        newview = WolfMapViewer(None, newcap, w=600, h=600)

        if self.linkedList is None:
            self.linkedList = [self]

        self.linkedList.append(newview)

        for curview in self.linkedList:
            curview.linked = True
            curview.linkedList = self.linkedList
            curview.link_shareopsvect = False

    def add_grid(self):
        mygrid = Grid(1000.)
        self.add_object('vector', newobj=mygrid, ToCheck=False, id='Grid')

    def add_WMS(self):
        xmin = 0
        xmax = 0
        ymin = 0
        ymax = 0
        orthos = {'IMAGERIE': {'1971': 'ORTHO_1971', '1994-2000': 'ORTHO_1994_2000',
                               '2006-2007': 'ORTHO_2006_2007',
                               '2009-2010': 'ORTHO_2009_2010',
                               '2012-2013': 'ORTHO_2012_2013',
                               '2015': 'ORTHO_2015', '2016': 'ORTHO_2016', '2017': 'ORTHO_2017',
                               '2018': 'ORTHO_2018', '2019': 'ORTHO_2019', '2020': 'ORTHO_2020',
                               '2021': 'ORTHO_2021'}}
        for idx, (k, item) in enumerate(orthos.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id='PPNC ' + m)
        self.add_object(which='wmsback',
                        newobj=imagetexture('PPNC', 'Orthos France', 'OI.OrthoimageCoverage.HR', '',
                                            self, xmin, xmax, ymin, ymax, -99999, 1024, France=True, epsg='EPSG:27563'),
                        ToCheck=False, id='Orthos France')

        forelist = {'EAU': {'Aqualim': 'RES_LIMNI_DGARNE', 'Alea': 'ALEA_INOND', 'Lidaxes': 'LIDAXES'},
                    'LIMITES': {'Secteurs Statistiques': 'LIMITES_QS_STATBEL'},
                    'INSPIRE': {'Limites administratives': 'AU_wms'},
                    'PLAN_REGLEMENT': {'Plan Percellaire': 'CADMAP_2021_PARCELLES'}}

        for idx, (k, item) in enumerate(forelist.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id=m)

    def set_compare(self, ListArrays=None):
        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison', w=600, h=600)
        third = WolfMapViewer(None, 'Difference', w=600, h=600)

        second.add_grid()
        third.add_grid()
        second.add_WMS()
        third.add_WMS()

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked = True
            curlist.linkedList = list

        if ListArrays is not None:
            if len(ListArrays) == 2:
                mnt = ListArrays[0]
                mns = ListArrays[1]
            else:
                return
        else:
            return

        mns: WolfArray
        mnt: WolfArray
        diff: WolfArray

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = mns - mnt

        mns.copy_mask(mnt, True)
        diff.copy_mask(mnt, True)

        mnt.change_gui(first)
        mns.change_gui(second)
        diff.change_gui(third)

        path = os.path.dirname(__file__)
        fn = join(path, 'models\\diff16.pal')
        diff.mypal.readfile(fn)
        diff.mypal.automatic = False
        diff.myops.palauto.SetValue(0)

        mnt.mypal.automatic = False
        mns.mypal.automatic = False
        mnt.myops.palauto.SetValue(0)
        mns.myops.palauto.SetValue(0)

        mnt.mypal.isopop(mnt.array, mnt.nbnotnull)
        mns.mypal.updatefrompalette(mnt.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=mnt, ToCheck=True, id='source')
        second.add_object('array', newobj=mns, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')

        mnt.myops.myzones = mns.myops.myzones
        diff.myops.myzones = mns.myops.myzones

        first.active_array = mnt
        second.active_array = mns
        third.active_array = diff

    def set_compare_all(self, ListArrays=None):
        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison MNS', w=600, h=600)
        third = WolfMapViewer(None, 'Difference MNS', w=600, h=600)
        if len(ListArrays) == 3:
            fourth = WolfMapViewer(None, 'Comparison MNT', w=600, h=600)
            fifth = WolfMapViewer(None, 'Difference MNT', w=600, h=600)

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        if len(ListArrays) == 3:
            list.append(fourth)
            list.append(fifth)

        for curview in list:
            if curview is not self:
                curview.add_grid()
                curview.add_grid()

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curview in list:
            curview.linked = True
            curview.linkedList = list

        comp2 = None
        if ListArrays is not None:
            if len(ListArrays) == 2:
                src = ListArrays[0]
                comp1 = ListArrays[1]
            elif len(ListArrays) == 3:
                src = ListArrays[0]
                comp1 = ListArrays[1]
                comp2 = ListArrays[2]
            else:
                return
        else:
            return

        src: WolfArray
        comp1: WolfArray
        diff1: WolfArray
        comp2: WolfArray
        diff2: WolfArray

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff1 = comp1 - src

        comp1.copy_mask(src, True)
        diff1.copy_mask(src, True)

        src.change_gui(first)
        comp1.change_gui(second)
        diff1.change_gui(third)

        src.mypal.automatic = False
        comp1.mypal.automatic = False
        src.myops.palauto.SetValue(0)
        comp1.myops.palauto.SetValue(0)

        src.mypal.isopop(src.array, src.nbnotnull)
        comp1.mypal.updatefrompalette(src.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=src, ToCheck=True, id='source')
        second.add_object('array', newobj=comp1, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff1, ToCheck=True, id='diff=comp-source')

        comp1.myops.myzones = src.myops.myzones
        diff1.myops.myzones = src.myops.myzones

        first.active_array = src
        second.active_array = comp1
        third.active_array = diff1

        if comp2 is not None:
            diff2 = comp2 - src
            comp2.copy_mask(src, True)
            diff2.copy_mask(src, True)

            comp2.change_gui(fourth)
            diff2.change_gui(fifth)

            comp2.mypal.automatic = False
            comp2.myops.palauto.SetValue(0)

            comp2.mypal.updatefrompalette(src.mypal)

            # Ajout des matrices dans les fenêtres de visualisation
            fourth.add_object('array', newobj=comp2, ToCheck=True, id='comp2')
            fifth.add_object('array', newobj=diff2, ToCheck=True, id='diff2=comp2-source')

            comp2.myops.myzones = src.myops.myzones
            diff2.myops.myzones = src.myops.myzones

            fourth.active_array = comp2
            fifth.active_array = diff2

    def set_blender_sculpting(self):

        myframe = wx.Frame(None, title=_('Excavation and backfill'))
        sizergen = wx.BoxSizer(wx.VERTICAL)
        sizer1 = wx.BoxSizer(wx.HORIZONTAL)
        sizer2 = wx.BoxSizer(wx.HORIZONTAL)
        sizer3 = wx.BoxSizer(wx.HORIZONTAL)
        sizergen.Add(sizer1)
        sizergen.Add(sizer2)
        sizergen.Add(sizer3)

        labexc = wx.StaticText(myframe, label=_('Excavation : '))
        labback = wx.StaticText(myframe, label=_('Backfill   : '))
        labbal = wx.StaticText(myframe, label=_('Balance   : '))
        sizer1.Add(labexc)
        sizer2.Add(labback)
        sizer3.Add(labbal)

        font = wx.Font(18, wx.DECORATIVE, wx.NORMAL, wx.NORMAL)

        Exc = wx.StaticText(myframe, label=' [m³]')
        Back = wx.StaticText(myframe, label=' [m³]')
        Bal = wx.StaticText(myframe, label=' [m³]')

        labexc.SetFont(font)
        labback.SetFont(font)
        labbal.SetFont(font)
        Exc.SetFont(font)
        Back.SetFont(font)
        Bal.SetFont(font)

        sizer1.Add(Exc)
        sizer2.Add(Back)
        sizer3.Add(Bal)

        myframe.SetSizer(sizergen)
        myframe.Layout()
        myframe.Centre(wx.BOTH)
        myframe.Show()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['ExcavationBackfill'] = myframe
        self.link_params['Excavation'] = Exc
        self.link_params['Backfill'] = Back
        self.link_params['Balance'] = Bal

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Sculpting', w=600, h=600)
        third = WolfMapViewer(None, 'Difference', w=600, h=600)
        fourth = WolfMapViewer(None, 'Gradient', w=600, h=600)
        fifth = WolfMapViewer(None, 'Laplace', w=600, h=600)
        sixth = WolfMapViewer(None, 'Unitary Mask', w=600, h=600)

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        list.append(fourth)
        list.append(fifth)
        list.append(sixth)

        for curlist in list:
            curlist.add_grid()
            curlist.add_WMS()

            # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked = True
            curlist.linkedList = list

        source: WolfArray
        sourcenew: WolfArray
        diff: WolfArray
        grad: WolfArray
        lap: WolfArray
        unimask: WolfArray

        source = self.active_array
        sourcenew = WolfArray(mold=source)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = source - source
        grad = source.get_gradient_norm()
        lap = source.get_laplace()
        unimask = WolfArray(mold=diff)

        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

        grad.copy_mask(source, True)
        lap.copy_mask(source, True)
        diff.copy_mask(source, True)
        unimask.copy_mask(source, True)

        sourcenew.change_gui(second)
        diff.change_gui(third)
        grad.change_gui(fourth)
        lap.pchange_gui(fifth)
        unimask.change_gui(sixth)

        path = os.path.dirname(__file__)
        # fn=join(path,'models\\diff16.pal')
        # diff.mypal.readfile(fn)
        # diff.mypal.automatic=False
        # diff.myops.palauto.SetValue(0)

        # fn=join(path,'models\\diff3.pal')
        # unimask.mypal.readfile(fn)
        # unimask.mypal.automatic=False
        # unimask.myops.palauto.SetValue(0)

        # Ajout des matrices dans les fenêtres de visualisation
        second.add_object('array', newobj=sourcenew, ToCheck=True, id='source_new')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')
        fourth.add_object('array', newobj=grad, ToCheck=True, id='gradient')
        fifth.add_object('array', newobj=lap, ToCheck=True, id='laplace')
        sixth.add_object('array', newobj=unimask, ToCheck=True, id='unimask')

        sourcenew.myops.myzones = source.myops.myzones
        diff.myops.myzones = source.myops.myzones
        grad.myops.myzones = source.myops.myzones
        lap.myops.myzones = source.myops.myzones
        unimask.myops.myzones = source.myops.myzones

        second.active_array = sourcenew
        third.active_array = diff
        fourth.active_array = grad
        fifth.active_array = lap
        sixth.active_array = unimask

        self.mimicme()

    def update_blender_sculpting(self):

        if not self.linked:
            return
        if len(self.linkedList) != 6:
            return

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self.linkedList[0]
        second = self.linkedList[1]
        third = self.linkedList[2]
        fourth = self.linkedList[3]
        fifth = self.linkedList[4]
        sixth = self.linkedList[5]

        source = first.active_array
        sourcenew = second.active_array
        diff = third.active_array
        grad = fourth.active_array
        lap = fifth.active_array
        unimask = sixth.active_array

        fn = ''
        if self.link_params is not None:
            if 'gltf file' in self.link_params.keys():
                fn = self.link_params['gltf file']
                fnpos = self.link_params['gltf pos']

        if fn == '':
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = self.link_params['gltf file']
                        fnpos = self.link_params['gltf pos']
                        break

        sourcenew.import_from_gltf(fn, fnpos)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff.array = (sourcenew - source).array
        grad.array = sourcenew.get_gradient_norm().array
        lap.array = sourcenew.get_laplace().array
        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

        diff.copy_mask(sourcenew, True)
        lap.copy_mask(sourcenew, True)
        grad.copy_mask(sourcenew, True)
        unimask.copy_mask(sourcenew, True)

        first.Paint()
        second.Paint()
        third.Paint()
        fourth.Paint()
        fifth.Paint()
        sixth.Paint()

        Exc: wx.StaticText
        Back: wx.StaticText
        Bal: wx.StaticText
        if not 'ExcavationBackfill' in self.link_params.keys():
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'ExcavationBackfill' in curgui.link_params.keys():
                        myframe = curgui.link_params['ExcavationBackfill']
                        Exc = curgui.link_params['Excavation']
                        Back = curgui.link_params['Backfill']
                        Bal = curgui.link_params['Balance']
        else:
            myframe = self.link_params['ExcavationBackfill']
            Exc = self.link_params['Excavation']
            Back = self.link_params['Backfill']
            Bal = self.link_params['Balance']

        Exc.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array < 0.])) + ' [m³]')
        Back.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array > 0.])) + ' [m³]')
        Bal.SetLabel("{:.2f}".format(np.sum(diff.array)) + ' [m³]')

    def zoomon_activevector(self, size=500., forceupdate=True):

        curvec = self.active_vector
        if curvec.minx == -99999:
            curvec.find_minmax()

        bounds = [curvec.minx, curvec.maxx, curvec.miny, curvec.maxy]

        dx = bounds[1] - bounds[0]
        dy = bounds[3] - bounds[2]

        self.mousex = bounds[0] + dx / 2.
        self.mousey = bounds[2] + dy / 2.
        self.width = max(size, dx)
        self.height = max(size, dy)

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def read_project(self, fn):
        myproject = Wolf_Param(None, filename=fn, toShow=False)

        mykeys = ['cross_sections', 'vector', 'array']

        if 'which' in myproject.myparams.keys():
            which = myproject.myparams['which']['action']['value']
            if which == 'compare':
                ListCompare = []
                if 'array' in myproject.myparams.keys():
                    for curid, curname in zip(myproject.myparams['array'].keys(), myproject.myparams['array'].values()):
                        ListCompare.append(WolfArray(normpath(curname['value'])))

                self.set_compare(ListCompare)
                return

        if 'cross_sections' in myproject.myparams.keys():
            for curid, curname in zip(myproject.myparams['cross_sections'].keys(),
                                      myproject.myparams['cross_sections'].values()):
                if curid != 'format' and curid != 'dirlaz':
                    mycs = crosssections(curname['value'],
                                         format=myproject.myparams['cross_sections']['format']['value'],
                                         dirlaz=myproject.myparams['cross_sections']['dirlaz']['value'])

                    self.add_object('cross_sections', newobj=mycs, id=curid)

        if 'vector' in myproject.myparams.keys():
            for curid, curname in zip(myproject.myparams['vector'].keys(), myproject.myparams['vector'].values()):
                myvec = Zones(curname['value'])
                self.add_object('vector', newobj=myvec, id=curid)

        if 'array' in myproject.myparams.keys():
            for curid, curname in zip(myproject.myparams['array'].keys(), myproject.myparams['array'].values()):
                curarray = WolfArray(curname['value'])
                self.add_object('array', newobj=curarray, id=curid)

        if 'wolf2d' in myproject.myparams.keys():
            for curid, curname in zip(myproject.myparams['wolf2d'].keys(), myproject.myparams['wolf2d'].values()):
                curwolf = Wolfresults_2D(curname['value'])
                self.add_object('res2d', newobj=curwolf, id=curid)

            self.menu_wolf2d()

        if 'palette' in myproject.myparams.keys():
            self.project_pal = {}
            for curid, curname in zip(myproject.myparams['palette'].keys(), myproject.myparams['palette'].values()):
                mypal = wolfpalette(None, '')
                mypal.readfile(curname['value'])
                mypal.automatic = False

                self.project_pal[curid] = mypal

        if 'palette-array' in myproject.myparams.keys():
            curarray: WolfArray
            if self.project_pal is not None:
                for curid, curname in zip(myproject.myparams['palette-array'].keys(),
                                          myproject.myparams['palette-array'].values()):
                    if curname['value'] in self.project_pal.keys():
                        curarray = self.getobj(curid)
                        if curarray is not None:
                            mypal = self.project_pal[curname['value']]
                            curarray.mypal = mypal
                            curarray.myops.palauto.SetValue(0)
                            curarray.updatepalette(0)
                            curarray.delete_lists()

        if 'cross_sections_link' in myproject.myparams.keys():
            if 'linkzones' in myproject.myparams['cross_sections_link'].keys():
                idx = myproject.myparams['cross_sections_link']['linkzones']['value']

                for curvect in self.added['vectors']:
                    myzones: Zones
                    myzones = self.added['vectors'][curvect]['values']
                    if myzones.idx == idx:
                        self.active_cs.link_external_zones(myzones)

                zonename = ''
                vecname = ''

                if 'sortzone' in myproject.myparams['cross_sections_link'].keys():
                    zonename = myproject.myparams['cross_sections_link']['sortzone']['value']
                if 'sortname' in myproject.myparams['cross_sections_link'].keys():
                    vecname = myproject.myparams['cross_sections_link']['sortname']['value']

                if zonename != '' and vecname != '':
                    names = [cur.myname for cur in myzones.myzones]
                    idx = names.index(zonename)
                    curzone = myzones.myzones[idx]
                    names = [cur.myname for cur in curzone.myvectors]
                    idx = names.index(vecname)
                    curvec = curzone.myvectors[idx]

                    if curvec is not None:
                        curvec: vector
                        self.active_cs.sort_along(curvec.asshapely_ls(), curvec.myname, False)

    def save_project(self, fn):
        myproject = Wolf_Param(None, toShow=False)

        mykeys = ['cross_sections', 'vector', 'array', 'wolf2d']
        for curkey in mykeys:
            myproject[curkey] = {}
        '''
        # myproject.myparams['which']={}
        # myproject.myparams['which']['action']={}
        # myproject.myparams['which']['action']['value']

        # mycs = self.active_cs
        # if mycs is not None:
        #     myproject.myparams['cross_sections']={}
        #     myproject.myparams['cross_sections']['mycs']={}
        #     myproject.myparams['cross_sections']['mycs']['value']=mycs.filename

        #     myproject.myparams['vector']={}
        #     myproject.myparams['vector']['river']={}
        #     myproject.myparams['vector']['river']['value']=self.added['vectors'][0].filename

        # if 'array' in myproject.myparams.key():
        #     for curid,curname in zip(myproject.myparams['array'].keys(),myproject.myparams['array'].values()):
        #         curarray=WolfArray(curname['value'])
        #         self.add_object('array',newobj=curarray,id=curid)
        '''
        # Dessin des matrices
        try:
            for curarray in self.added['arrays']:
                if self.added['arrays'][curarray]['checked']:
                    locarray: WolfArray
                    locarray = self.added['arrays'][curarray]['values']
                    myproject['array'][locarray.idx] = {}
                    myproject['array'][locarray.idx]['value'] = locarray.filename
        except:
            pass

        # Dessin des résultats 2D
        try:
            for curarray in self.added['wolf2d']:
                if self.added['wolf2d'][curarray]['checked']:
                    locarray: wolf2dprev
                    locarray = self.added['wolf2d'][curarray]['values']
                    myproject['wolf2d'][locarray.idx] = {}
                    myproject['wolf2d'][locarray.idx]['value'] = locarray.filename
        except:
            pass

        # Dessin des vecteurs
        try:
            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    curvec: Zones
                    curvec = self.added['vectors'][curvect]['values']
                    myproject['wolf2d'][curvec.idx] = {}
                    myproject['wolf2d'][curvec.idx]['value'] = curvec.filename
        except:
            pass

    def read_laz(self, fn=None):
        if fn is None:
            filternpz = "npz (*.npz)|*.npz|all (*.*)|*.*"
            dlg = wx.FileDialog(None, _('Choose LAS npz file'), wildcard=filternpz)
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            fn = dlg.GetPath()

        self.mylazdata = read_laz(fn)

        if self.linked:
            if len(self.linkedList) > 0:
                for curframe in self.linkedList:
                    curframe.mylazdata = self.mylazdata

    def managebanks(self):
        if self.notebookbanks is None:
            self.notebookbanks = PlotNotebook(self)
            self.mypagebanks = self.notebookbanks.add(_("Manager banks interpolator"), "ManagerInterp")

        msg = ''
        if self.active_cs is None:
            msg += _(' The is no cross section. Please activate the desired object !')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        if self.active_cs.linked_zones is None:
            msg += _(' The active zones is None. Please link the desired object to the cross sections !\n')
        # if self.active_zone is None:
        #     msg+=_(' The active zone is None. Please activate the desired object !\n')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        self.mypagebanks.pointing(self, self.active_cs, self.active_vector)
        self.notebookbanks.Show(True)

    def _set_fn_fnpos_gltf(self):
        dlg = wx.FileDialog(None, _('Choose filename'),
                            wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fn = dlg.GetPath()
        dlg.Destroy()

        dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                            style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fnpos = dlg.GetPath()
        dlg.Destroy()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['gltf file'] = fn
        self.link_params['gltf pos'] = fnpos

        return fn

    def set_fn_fnpos_gltf(self):
        fn = ''
        fnpos = ''
        if self.linked:
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = curgui.link_params['gltf file']
                        fnpos = curgui.link_params['gltf pos']
                        break
        elif self.link_params is None:
            self.link_params = {}
            fn = self._set_fn_fnpos_gltf()

        if fn == '':
            self._set_fn_fnpos_gltf()

    def OnMenubar(self, event: wx.CommandEvent):
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        autoscale = True

        if id == wx.ID_OPEN:
            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Choose file", wildcard=filterProject)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            os.chdir(os.path.dirname(filename))
            self.read_project(filename)

        elif itemlabel == _('Update .bloc file'):
            
            msg = _('If you continue the .bloc file will be relpaced !')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'
            
            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return
                        
            self.wolfparent.write_bloc_file()
            
        elif itemlabel == _('Update free surface elevation - IC'):
            
            if len(self.active_array.mngselection.myselection)==0:
                
                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'
                
                wx.LogWarning(msg)
                return                
                        
            self.wolfparent.extend_freesurface_elevation(self.active_array.mngselection.myselection)
            
        elif itemlabel== _('Update roughness coeff'):
            
            if len(self.active_array.mngselection.myselection)==0:
                
                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'
                
                wx.LogWarning(msg)
                return                
                        
            self.wolfparent.extend_roughness(self.active_array.mngselection.myselection)
            
        elif itemlabel == _('Update model from current mask'):

            if type(self.active_array) != WolfArray_Sim2D:
                msg = _('Please select a mono-block array !')+'\n'
                dlg=wx.MessageBox(msg,style=wx.OK)
                return

            msg = _('If you continue, the mask of all arrays will be replaced by the current mask !')+'\n'
            msg += _('The external contour in the .bloc file will also be relpaced.')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'
            
            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return
            
            with wx.lib.busy.BusyInfo(_('Updating 2D model')):
                wait = wx.BusyCursor()
            
                sux,suy,cont,interior = self.active_array.suxsuy_contour(self.wolfparent.filenamegen,True)            
            
                self.wolfparent.mimic_mask(self.active_array)
                self.wolfparent.replace_external_contour(cont,interior)
                
                del wait
                        
            self.wolfparent.extend_bed_elevation()

        elif itemlabel == _("Plot active vector..."):

            if self.active_vector is None:
                return

            fig, ax = self.active_vector.plot_mpl(True, False)

            linkedarrays, labels = self.get_linked_arrays()

            self.active_vector.plot_linked(fig, ax, linkedarrays, labels)

        elif itemlabel == _("Plot active polygons..."):

            if self.active_zone is None:
                return

            plotzone = []
            zonename = self.active_zone.myname
            if '_left_' in zonename or '_right_' in zonename:

                testname = zonename.replace('_left_', '')
                testname = testname.replace('_right_', '')

                for curzone in self.active_zones.myzones:
                    if testname == curzone.myname.replace('_left_', '').replace('_right_', ''):
                        plotzone.append(curzone)

                msg = wx.MessageDialog(self,
                                       _('Left and Right polygons are detected \nWould you like to plot left and right polygons in the same plot ?'),
                                       style=wx.YES_NO | wx.YES_DEFAULT)
                ret = msg.ShowModal()
                msg.Destroy()
                if ret == wx.ID_NO:
                    plotzone = [self.active_zone]
            else:
                plotzone = [self.active_zone]

            fig, ax = plt.subplots(1, 1)

            linkedarrays = []
            labels = []

            # Matrices 2D
            for curarray in self.added['arrays']:
                if self.added['arrays'][curarray]['checked']:
                    locarray = self.added['arrays'][curarray]['values']
                    linkedarrays.append(locarray)
                    labels.append(locarray.idx)

            # Résultats 2D
            for curarray in self.added['wolf2d']:
                if self.added['wolf2d'][curarray]['checked']:
                    locarray = self.added['wolf2d'][curarray]['values']
                    linkedarrays.append(locarray)
                    labels.append(locarray.idx)

            if len(plotzone) > 1:
                pass
                for curzone in plotzone:
                    if '_left_' in curzone.myname:
                        loclabels = [cur + '_left' for cur in labels]
                        curzone.plot_linked_poly(fig, ax, linkedarrays, loclabels, '--')
                    elif '_right_' in curzone.myname:
                        loclabels = [cur + '_right' for cur in labels]
                        curzone.plot_linked_poly(fig, ax, linkedarrays, loclabels, '-.')
            else:
                self.active_zone.plot_linked_poly(fig, ax, linkedarrays, labels)

            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    myzones: Zones
                    myzones = self.added['vectors'][curvect]['values']
                    if myzones.idx == 'bridges' or myzones.idx == 'ponts':

                        curzone: zone
                        curzone = myzones.myzones[0]

                        trace: vector
                        trace = curzone.myvectors[0]
                        tracels = trace.asshapely_ls()

                        for curvect in curzone.myvectors:
                            curvect: vector
                            if curvect is not trace:
                                curls = curvect.asshapely_ls()

                                if curls.intersects(tracels):
                                    inter = curls.intersection(tracels)
                                    curs = float(tracels.project(inter))
                                    locz = np.asarray([vert.z for vert in curvect.myvertices])
                                    zmin = np.amin(locz)
                                    zmax = np.amax(locz)

                                    ax.scatter(curs, zmin, label=curvect.myname + ' min')
                                    ax.scatter(curs, zmax, label=curvect.myname + ' max')

            ax.legend()
            fig.show()

        elif itemlabel == _("Change current view"):

            dlg = wx.SingleChoiceDialog(None, "Pick a view", "Choices", CHOICES_VIEW_2D)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            dlg = wx.MessageDialog(None, _('Apply to all results?'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:
                self.active_res2d.set_currentview(method)
            else:
                for curarray in self.added['wolf2d']:
                    if self.added['wolf2d'][curarray]['checked']:
                        locarray = self.added['wolf2d'][curarray]['values']
                        locarray.set_currentview(method)

        elif itemlabel == _("Read last result"):

            self.currently_readresults = True

            for curmodel in self.added['wolf2d']:
                if self.added['wolf2d'][curmodel]['checked']:
                    locmodel: Wolfresults_2D
                    locmodel = self.added['wolf2d'][curmodel]['values']
                    locmodel.read_oneresult()
                    locmodel.set_currentview(locmodel.currentview)

            self.currently_readresults = False

        elif itemlabel == _("Manage boundary conditions..."):
            if self.active_res2d is not None:
                self.active_res2d.myparams.editing_bc(self.added['wolf2d'])

        elif itemlabel == _("Manage banks..."):
            if self.active_vector is None:
                msg = _('Active vector is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.managebanks()

        elif itemlabel == _("Create banks from vertices..."):

            self.active_cs.create_zone_from_banksbed()
            self.active_cs.linked_zones.showstructure(self)

        elif itemlabel == _("Link cross sections to active zones"):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            if self.active_zones is None:
                msg = _('Active zone is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.active_cs.link_external_zones(self.active_zones)

        elif itemlabel == _("Rename cross sections..."):

            dlg = wx.TextEntryDialog(None, _('Which starting point?'))
            ret = dlg.ShowModal()

            idxstart = dlg.GetValue()

            self.active_cs.rename(int(idxstart))

        elif itemlabel == _("Triangulate cross sections..."):
            self.triangulate_cs()

        elif itemlabel == "Create bridge and export gltf...":

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.action = 'bridge gltf'

        elif itemlabel == _("Export cross sections to gltf..."):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            dlg = wx.TextEntryDialog(self, 'Z minimum ?', 'Choose an elevation as base')
            dlg.SetValue('')

            zmin = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmin = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            self.active_cs.export_gltf(zmin, fn)

        elif itemlabel == _("New cloud Viewer..."):
            if self.myinterp is not None:
                self.myinterp.viewer_interpolator()

        elif itemlabel == _("Interpolate on active array..."):
            if self.myinterp is not None:
                self.interpolate_cs()

        elif itemlabel == _('Save project'):
            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Name your file", wildcard=filterProject, style=wx.FD_SAVE)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            self.save_project(filename)
        elif itemlabel == _('Read LAZ data from npz'):
            self.read_laz()

        elif itemlabel == _('Create cloud points from bridges'):
            if self.mylazdata is None:
                self.read_laz()

            mybridges = self.mylazdata[np.where(self.mylazdata[:, 3] == 10)]
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybridges)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([255, 102, 102])
            mycloud.myprop.width = .5

            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')

        elif itemlabel == _('Create cloud points from buildings'):
            if self.mylazdata is None:
                self.read_laz()

            mybuildings = self.mylazdata[np.where(self.mylazdata[:, 3] == 1)]
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybuildings)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([102, 102, 102])
            mycloud.myprop.width = .5
            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')

        elif itemlabel == _('Create LAZ viewer'):
            if self.mylazdata is None:
                self.read_laz()
            myviewer(self.mylazdata, 0)

        elif itemlabel == _('Clip LAZ data on current zoom'):
            filternpz = "npz (*.npz)|*.npz|all (*.*)|*.*"
            dlg = wx.DirDialog(None, _('Choose XYZ Numpy LAZ Directory'))
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            dlg = wx.FileDialog(None, _('Choose LAS npz out file'), wildcard=filternpz, style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            fn = dlg.GetPath()

            self.mylazdir = dlg.GetPath()
            curbounds = [[self.xmin, self.xmin + self.width], [self.ymin, self.ymin + self.height]]
            clip_data_xyz(self.mylazdir, fn, curbounds)
            self.mydatalaz = read_laz(fn)

        elif itemlabel == _('Multiviewer'):
            dlg = wx.NumberEntryDialog(self, _("Additional viewers"), _("How many?"), _("How many additional viewers?"),
                                       1, 0, 5)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            nb = dlg.GetValue()
            dlg.Destroy()
            for i in range(nb):
                self.add_viewer_and_link()

        elif itemlabel == _('Set comparison'):
            ListArrays = []
            filterProject = "Array (*.bin)|*.bin|all (*.*)|*.*"
            file = wx.FileDialog(self, "Choose Source file", wildcard=filterProject)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                ListArrays.append(file.GetPath())
                file.Destroy()

            file = wx.FileDialog(self, _("Choose Comparison file"), wildcard=filterProject)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                ListArrays.append(file.GetPath())
                file.Destroy()

            first = WolfArray(ListArrays[0], preload=False)
            sec = WolfArray(ListArrays[1], preload=False)
            first.read_txt_header()
            sec.read_txt_header()

            if first.nbx == sec.nbx and first.nby == sec.nby and first.origx == sec.origx and first.origy == sec.origy:
                first.preload = True
                sec.preload = True
                first.read_data()
                sec.read_data()
                first.mask_data(0.)
                sec.mask_data(0.)
                self.set_compare([first, sec])

        elif id == wx.ID_FILE1:
            self.add_object(which='array', ToCheck=True)
            # self.mybc=BcManager(self,newarray.dx,newarray.dy,newarray.origx,newarray.origy,"Boundary conditions")
            # self.mybc.FindBorders(newarray.array)
        elif id == wx.ID_FILE2:
            self.add_object(which='vector', ToCheck=True)
        elif id == wx.ID_FILE3:
            self.add_object(which='cloud', ToCheck=True)
        elif id == wx.ID_FILE4:
            self.add_object(which='cross_sections', ToCheck=True)
        elif itemlabel == _('Add Wolf2D results...'):
            self.add_object(which='res2d', ToCheck=True)
            self.menu_wolf2d()

        elif itemlabel == _('Add array and crop...'):
            self.add_object(which='array_crop', ToCheck=True)

        elif itemlabel == _('Create array from bathymetry file...'):

            self.add_object(which='array_xyz', ToCheck=True)

        elif itemlabel == _('Create array from Lidar 2002...'):

            dlg = wx.SingleChoiceDialog(None, _('What source of data?'), _('Lidar 2002'),
                                        [_('First echo'), _('Second echo')])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                return

            sel = dlg.GetStringSelection()

            if sel == _('First echo'):
                self.add_object(which='array_lidar_first', ToCheck=True)
            elif sel == _('Second echo'):
                self.add_object(which='array_lidar_second', ToCheck=True)

        elif id == wx.ID_FILE5:
            def addscandir(mydir):
                for entry in scandir(mydir):
                    if entry.is_dir():
                        addscandir(entry)
                    elif entry.is_file():
                        if entry.name.endswith('.vec') or entry.name.endswith('.vecz'):

                            msg = wx.MessageDialog(self,
                                                   _(entry.name + ' found in ' + mydir + '\n\n Is it a "cross sections" file?'),
                                                   style=wx.YES_NO | wx.NO_DEFAULT)
                            ret = msg.ShowModal()
                            if ret == wx.ID_YES:
                                self.add_object(which='vector',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))
                            else:
                                self.add_object(which='cross_sections',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))

                        elif entry.name.endswith('.bin'):
                            self.add_object(which='array',
                                            filename=join(mydir, entry.name),
                                            ToCheck=True,
                                            id=join(mydir, entry.name))

            mydialog = wx.DirDialog(self, _("Choose directory to scan"))
            if mydialog.ShowModal() == wx.ID_CANCEL:
                mydialog.Destroy()
                return
            else:
                # récupération du nom de fichier avec chemin d'accès
                mydir = mydialog.GetPath()
                mydialog.Destroy()

            if exists(mydir):
                addscandir(mydir)

        elif id == wx.ID_FILE6:
            # Création d'une nouvelle matrice
            newarray = WolfArray(create=True, parentgui=self)
            self.add_object('array', newobj=newarray)
        elif id == wx.ID_FILE7:
            # Création de nouveaux vecteurs
            newzones = Zones(parent=self)
            self.add_object('vector', newobj=newzones)
        elif id == wx.ID_FILE8:
            # Création d'un nouveau nuage de point
            newcloud = cloud_vertices()
            self.add_object('cloud', newobj=newcloud)
        elif id in self.tools.keys():
            # gestion des actions
            self.ManageActions(id)
        elif id == wx.ID_SAVE:

            for curarray in self.added['arrays']:
                if self.added['arrays'][curarray]['checked']:
                    obj: WolfArray
                    obj = self.added['arrays'][curarray]['values']

                    if obj.filename == '':
                        filterArray = "bin (*.bin)|*.bin|all (*.*)|*.*"
                        fdlg = wx.FileDialog(self, "Choose file", wildcard=filterArray, style=wx.FD_SAVE)
                        fdlg.ShowModal()
                        if fdlg.ShowModal() == wx.ID_OK:
                            obj.filename = fdlg.GetPath()

                    obj.write_all()

            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']
                    obj.saveas()

        elif itemlabel == 'Save to image...':
            autoscale = False
            fn, ds = self.save_canvasogl()
            self.save_linked_canvas(fn[:-4], ds)
        elif itemlabel == _('Export to gltf...'):

            curarray: WolfArray
            curvec: vector

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')
            if self.active_vector is None:
                msg += _('Active vector is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array
            curvec = self.active_vector

            curvec.find_minmax()

            i1, j1 = curarray.get_ij_from_xy(curvec.minx, curvec.miny)
            x1, y1 = curarray.get_xy_from_ij(i1, j1)
            x1 -= curarray.dx / 2.
            y1 -= curarray.dy / 2.

            i2, j2 = curarray.get_ij_from_xy(curvec.maxx, curvec.maxy)
            x2, y2 = curarray.get_xy_from_ij(i2, j2)
            x2 += curarray.dx / 2.
            y2 += curarray.dy / 2.
            mybounds = [[x1, x2], [y1, y2]]

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            curarray.export_to_gltf(mybounds, fn)

        elif itemlabel == _('Import from gltf...'):

            curarray: WolfArray

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                                style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fnpos = dlg.GetPath()
            dlg.Destroy()

            curarray.import_from_gltf(fn, fnpos)

        elif itemlabel == _('Set gltf comparison'):

            self.set_blender_sculpting()
            autoscale = False

        elif itemlabel == _('Update from gltf...'):
            autoscale = False

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_fn_fnpos_gltf()
            self.update_blender_sculpting()

        elif id == wx.ID_SAVEAS:

            for curarray in self.added['arrays']:
                if self.added['arrays'][curarray]['checked']:
                    obj: WolfArray
                    obj = self.added['arrays'][curarray]['values']

                    filterArray = "bin (*.bin)|*.bin|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Array : " + obj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.filename = fdlg.GetPath()
                        obj.write_all()

            for curvect in self.added['vectors']:
                if curvect == 'grid':
                    pass
                elif self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']

                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + obj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.saveas(fdlg.GetPath())

        if len(self.myarrays) + len(self.myvectors) + len(self.myclouds) + len(self.myres2D) == 2 and autoscale:
            # Trouve les bornzs si un seul élément est présent, sinon on conserve l'état du zoom
            self.Autoscale()

    def get_linked_arrays(self):
        linkedarrays = []
        labels = []

        for curarray in self.added['arrays']:
            if self.added['arrays'][curarray]['checked']:
                locarray = self.added['arrays'][curarray]['values']
                linkedarrays.append(locarray)
                labels.append(locarray.idx)

        for curarray in self.added['wolf2d']:
            if self.added['wolf2d'][curarray]['checked']:
                locarray = self.added['wolf2d'][curarray]['values']
                linkedarrays.append(locarray)
                labels.append(locarray.idx)
        
        return linkedarrays,labels

    def save_linked_canvas(self, fn, mpl=True, ds=0.):
        if self.linked:
            for idx, curel in enumerate(self.linkedList):
                curel.save_canvasogl(fn + '_' + str(idx) + '.png', mpl, ds)

    def thread_update_blender(self):
        print("Update blender")
        if self.canvas.SetCurrent(self.context):
            self.update_blender_sculpting()
            t = threading.Timer(10.0, self.thread_update_blender)
            t.start()

    def add_object(self, which: str = 'array', filename='', newobj=None, ToCheck=True, id=''):
        filterArray = "bin (*.bin)|*.bin|Float ESRI (*.flt)|*.flt|Elevation WOLF2D (*.top)|*.top|all (*.*)|*.*"
        filterres2d = "all (*.*)|*.*"
        filterVector = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|all (*.*)|*.*"
        filterCloud = "xyz (*.xyz)|*.xyz|text (*.txt)|*.txt|all (*.*)|*.*"
        filterCs = "txt 2022 (*.txt)|*.txt|vecz WOLF (*.vecz)|*.vecz|WOLF (*.sxy)|*.sxy|text 2000 (*.txt)|*.txt|all (*.*)|*.*"

        if filename == '' and newobj is None:
            # ouverture d'une boîte de dialogue
            if which.lower() == 'array' or which.lower() == 'array_crop':
                file = wx.FileDialog(self, "Choose file", wildcard=filterArray)
            elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':
                file = wx.DirDialog(self, "Choose directory containing Lidar data")
            elif which.lower() == 'array_xyz':
                file = wx.DirDialog(self, "Choose directory containing XYZ files")
            elif which.lower() == 'vector':
                file = wx.FileDialog(self, "Choose file", wildcard=filterVector)
            elif which.lower() == 'cloud':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)
            elif which.lower() == 'cross_sections':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCs)
            elif which.lower() == 'other':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)
            elif which.lower() == 'res2d':
                file = wx.FileDialog(self, "Choose file", wildcard=filterres2d)
            elif which.lower() == 'wmsback':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)
            elif which.lower() == 'wmsfore':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                try:
                    curfilter = file.GetFilterIndex()
                except:
                    pass
                file.Destroy()

        if filename != '':
            if (not (os.path.exists(filename))):
                print("Warning : the following file is not present here : " + filename)
                return

        if which.lower() == 'array' or which.lower() == 'array_crop':
            curdict = self.added['arrays']
            curtree = self.myitemsarray

            if newobj is None:

                testobj = WolfArray()
                testobj.filename = filename
                testobj.read_txt_header()

                if testobj.wolftype in WOLF_ARRAY_MB:
                    newobj = WolfArrayMB(filename, parentgui=self)
                else:
                    if which.lower() == 'array_crop':
                        newobj = WolfArray(filename, parentgui=self, crop='newcrop')
                    else:
                        newobj = WolfArray(filename, parentgui=self)

            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            newobj.change_gui(self)
            self.active_array = newobj

        elif which.lower() == 'array_xyz':

            curdict = self.added['arrays']
            curtree = self.myitemsarray

            msg = wx.MessageDialog(self, _('Would you like to crop the data?'), style=wx.YES_NO | wx.YES_DEFAULT)
            ret = msg.ShowModal()
            msg.Destroy()

            if ret == wx.ID_YES:

                newcrop = CropDialog(None)

                badvalues = True
                while badvalues:
                    badvalues = False

                    ret = newcrop.ShowModal()
                    if ret == wx.ID_CANCEL:
                        newcrop.Destroy()
                        return
                    else:
                        cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                                   [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                        tmpdx = float(newcrop.dx.Value)
                        tmpdy = float(newcrop.dy.Value)

                newcrop.Destroy()

                myxyz = xyz_scandir(filename, cropini)

                myhead = newcrop.get_header()
                if min(myhead.dx, myhead.dy) != 1.:
                    myhead.nbx = int(myhead.nbx * myhead.dx)
                    myhead.nby = int(myhead.nby * myhead.dy)
                    myhead.dx = 1.
                    myhead.dy = 1.

            else:
                myxyz = xyz_scandir(filename, None)
                myhead = header_wolf()

                myhead.origx = np.min(myxyz[:, 0]) - .5
                myhead.origy = np.min(myxyz[:, 1]) - .5

                myhead.dx = 1.
                myhead.dy = 1.
                tmpdx = 1.
                tmpdy = 1.

                myhead.nbx = int(np.max(myxyz[:, 0]) - myhead.origx) + 1
                myhead.nby = int(np.max(myxyz[:, 1]) - myhead.origy) + 1

            if len(myxyz) == 0:
                return

            newobj = WolfArray()

            newobj.init_from_header(myhead)
            newobj.nullvalue = -99999.
            newobj.array.data[:, :] = -99999.

            newobj.fillin_from_xyz(myxyz)

            newobj.mask_data(newobj.nullvalue)

            if min(tmpdx, tmpdy) != 1.:
                newobj.rebin(min(tmpdx, tmpdy))
                newobj.mask_data(newobj.nullvalue)

            newobj.change_gui(self)
            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            self.active_array = newobj

        elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':
            curdict = self.added['arrays']
            curtree = self.myitemsarray

            newcrop = CropDialog(None)

            badvalues = True
            while badvalues:
                badvalues = False

                ret = newcrop.ShowModal()
                if ret == wx.ID_CANCEL:
                    newcrop.Destroy()
                    return
                else:
                    cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                               [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                    tmpdx = float(newcrop.dx.Value)
                    tmpdy = float(newcrop.dy.Value)

            newcrop.Destroy()

            first, sec = Lidar2002.lidar_scandir(filename, cropini)

            if which.lower() == 'array_lidar_first':
                if len(first) == 0:
                    return

                newobj = Lidar2002.create_wolfarray(first, bounds=cropini)

                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj

                id = 'lidar2002_firstecho'
            else:
                if len(sec) == 0:
                    return
                newobj = Lidar2002.create_wolfarray(sec, bounds=cropini)
                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj
                id = 'lidar2002_secondecho'

        elif which.lower() == 'res2d':
            curdict = self.added['wolf2d']
            curtree = self.myitemsres2d

            if newobj is None:
                newobj = Wolfresults_2D(filename)
                newobj.read_param_simul()

            newobj.updatepalette()
            self.myres2D.append(newobj)
            newobj.parentgui = self
            self.active_res2d = newobj

        elif which.lower() == 'vector':
            curdict = self.added['vectors']
            curtree = self.myitemsvector
            if newobj is None:
                newobj = Zones(filename)
            self.myvectors.append(newobj)

        elif which.lower() == 'cross_sections':
            curdict = self.added['vectors']
            curtree = self.myitemsvector

            if newobj is None:

                dlg = wx.MessageDialog(None, 'Load LAZ data?', style=wx.YES_NO | wx.NO_DEFAULT)
                ret = dlg.ShowModal()
                dlg.Destroy()
                dirlaz = ''

                if ret == wx.ID_YES:
                    dlg = wx.DirDialog(None, 'If exist, where are the LAZ data?')
                    ret = dlg.ShowModal()
                    if ret == wx.ID_OK:
                        dirlaz = dlg.GetPath()

                if curfilter == 0:  # txt 2022
                    newobj = crosssections(filename, format='2022', dirlaz=dirlaz)
                if curfilter == 1:  # vecz
                    newobj = crosssections(filename, format='vecz', dirlaz=dirlaz)
                elif curfilter == 2:  # sxy
                    newobj = crosssections(filename, format='sxy', dirlaz=dirlaz)
                else:  # txt 2000
                    newobj = crosssections(filename, format='2000', dirlaz=dirlaz)

            self.myvectors.append(newobj)
            newobj.parentgui = self

        elif which.lower() == 'cloud':
            curdict = self.added['clouds']
            curtree = self.myitemscloud
            if newobj is None:
                newobj = cloud_vertices(filename)
            self.myclouds.append(newobj)

        elif which.lower() == 'other':
            if not newobj is None:
                curdict = self.added['others']
                curtree = self.myitemsothers
                self.myothers.append(newobj)

        elif which.lower() == 'wmsback':
            if not newobj is None:
                curdict = self.added['wms-background']
                curtree = self.myitemswmsback
                self.mywmsback.append(newobj)

        elif which.lower() == 'wmsfore':
            if not newobj is None:
                curdict = self.added['wms-foreground']
                curtree = self.myitemswmsfore
                self.mywmsfore.append(newobj)

        if id == '':
            dlg = wx.TextEntryDialog(self, 'ID ? (case insensitive)', 'Choose an identifier', _('NewObject'))
            dlg.SetValue('')
            if len(curdict) == 0:
                if dlg.ShowModal() == wx.ID_OK:
                    id = dlg.GetValue()
            else:
                id = list(curdict.keys())[0]
                while id.lower() in curdict:
                    if dlg.ShowModal() == wx.ID_OK:
                        id = dlg.GetValue()
            dlg.Destroy()

        if id.lower() in curdict:
            endid = '_'
            while id.lower() in curdict:
                id += endid

        newobj.idx = id.lower()

        myitem = self.treelist.AppendItem(curtree, id, data=newobj)

        if ToCheck:
            self.treelist.CheckItem(myitem)
            self.treelist.CheckItem(self.treelist.GetItemParent(myitem))

            newobj.check_plot()

        loc = curdict[id.lower()] = {}
        if filename != '':
            loc['listname'] = filename.lower()
            loc['filename'] = filename.lower()
        loc['values'] = newobj
        loc['listitem'] = myitem
        loc['checked'] = ToCheck

        if type(newobj) == crosssections:
            self.myclouds.append(newobj.cloud)
            newobj.cloud.idx = newobj.idx
            newobj.cloud.myprop.filled = True
            myitem = self.treelist.AppendItem(self.myitemscloud, id, data=newobj.cloud)

            curdict = self.added['clouds']
            loc = curdict[id.lower()] = {}
            loc['values'] = newobj.cloud
            loc['listitem'] = myitem
            loc['checked'] = False

            for curarray in self.added['arrays']:
                if self.added['arrays'][curarray]['checked']:
                    newobj.linked_arrays.append(self.added['arrays'][curarray]['values'])
                    newobj.linked_labels.append(self.added['arrays'][curarray]['values'].idx)

        elif type(newobj) == WolfArray:
            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

            if self.active_cs is not None:
                mycs = self.active_cs
                if mycs is not None:
                    mycs.linked_arrays.append(newobj)
                    mycs.linked_labels.append(id)

    def getobj(self, id: str):
        for curdict in self.added:
            if id.lower() in self.added[curdict].keys():
                try:
                    return self.added[curdict][id]['values']
                except:
                    return None

    def OnShowPopup(self, event):
        pos = event.GetPosition()
        if pos == (-1, -1):
            width, height = self.GetSize()
            pos = (width / 2., height / 2.)
        # else:
        #     pos = pos - self.GetPosition()
        self.PopupMenu(self.popupmenu, pos)

    def OnPopupItemSelected(self, event):
        item = self.popupmenu.FindItemById(event.GetId())
        text = item.ItemLabel

        if text == 'Save':
            if self.selobj is not None:
                if type(self.selobj) is WolfArray or type(self.selobj) is WolfArray_Sim2D:
                    self.selobj.write_all()
                elif type(self.selobj) is Zones:
                    self.selobj.saveas()

        elif text == 'Save as':
            if self.selobj is not None:
                if type(self.selobj) is WolfArray or type(self.selobj) is WolfArray_Sim2D:
                    filterArray = "bin (*.bin)|*.bin|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Array : " + self.selobj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selobj.filename = fdlg.GetPath()
                        self.selobj.write_all()
                elif type(self.selobj) is Zones:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + self.selobj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selobj.saveas(fdlg.GetPath())

                        # wx.MessageBox("You selected item '%s'" % text)

    def OnClose(self, event):
        if self.linked:
            if self.linkedList is not None:
                if self in self.linkedList:
                    id = self.linkedList.index(self)
                    self.linkedList.pop(id)

        self.Destroy()

    # def OnActivatedItem(self,e):
    #     myitem=e.GetItem()
    #     nameitem=self.treelist.GetItemText(myitem).lower()
    #     self.selobj = self.getobj(nameitem)
    #     self.selitem.SetLabel(nameitem)
    #     pass

    def OnCheckItem(self, event):

        myitem = event.GetItem()
        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        if nameparent != '':
            curobj = self.getobj(nameitem)
            self.added[nameparent][nameitem]['checked'] = bool(check)

            if bool(check):
                curobj.check_plot()
            else:
                if type(curobj) in [WolfArray, WolfArrayMB]:
                    curobj.uncheck_plot(True)
                else:
                    curobj.uncheck_plot()

            # if nameparent == 'vectors' or nameparent == 'cross_sections':
            #     if wx.GetKeyState(wx.WXK_CONTROL):
            #         curobj.showstructure(self)
            
            if curobj.idx == 'grid' and check:
                dlg = wx.TextEntryDialog(self, 'Size of the Grid ? (float)', 'Choose an size')
                dlg.SetValue('1000.')
                size = 1000.
                if dlg.ShowModal() == wx.ID_OK:
                    size = float(dlg.GetValue())
                curobj.creategrid(size, self.xmin, self.ymin, self.xmax, self.ymax)

    def getXY(self, pospix):

        width, height = self.canvas.GetSize()
        X = float(pospix[0]) / self.sx + self.xmin
        Y = float(height - pospix[1]) / self.sy + self.ymin
        return X, Y

    def OnZoomGesture(self, e):
        pass

    def OnLeave(self, e):
        if e.ControlDown():
            self.mytooltip.Show(False)

    def get_cross_sections(self):
        for curvect in self.added['vectors']:
            if self.added['vectors'][curvect]['checked']:
                obj = self.added['vectors'][curvect]['values']
                if type(obj) is crosssections:
                    return obj

        return None

    def OnRightDown(self, e: wx.MouseEvent):
        pos = e.GetPosition()
        x, y = self.getXY(pos)

        alt = e.AltDown()
        ctrl = e.ControlDown()
        shiftdown = e.ShiftDown()

        if self.action is None:
            self.rightdown = (x, y)

        elif self.action == 'bridge gltf':

            self.bridgepar = (x, y)

            dlg = wx.TextEntryDialog(self, 'Z maximum ?', 'Choose an elevation as top')
            dlg.SetValue('')

            zmax = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmax = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            points, triangles = self.active_vector.triangulation_ponts(self.bridgepar[0], self.bridgepar[1], zmax)
            self.active_cs.export_gltf_gen(points, triangles, fn)

            self.action = ''

        elif self.action == 'Select nearest profile':

            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")

            curprofile: profile

            for curvect in self.added['vectors']:
                if self.added['vectors'][curvect]['checked']:
                    obj = self.added['vectors'][curvect]['values']
                    if type(obj) is crosssections:
                        if self.myfigcs.mycs is not None:
                            self.myfigcs.mycs.myprop.width = 1
                            self.myfigcs.mycs.myprop.color = 0

                        curprofile = obj.select_profile(x, y)
                        self.myfigcs.set_cs(curprofile)
                        self.myfigcs.plot_cs()

                        curprofile.myprop.width = 2
                        curprofile.myprop.color = getIfromRGB([255, 0, 0])

                        self.Paint()
        elif self.action.find('select active vector') > -1:

            inside = self.action.find('inside') > -1
            onlyonezone = self.action.find('2') > -1

            if onlyonezone:
                self.active_zone.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zone.get_selected_vectors()

                if self.active_vector is not None:
                    self.active_zone.parent.Activate_vector(self.active_vector)
                    self.active_zone.active_vector = self.active_vector
                    self.active_zones.active_zone = self.active_vector.parentzone
            else:
                self.active_zones.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zones.get_selected_vectors()

            if self.active_vector is not None:
                self.active_zones.Activate_vector(self.active_vector)
                self.active_zone = self.active_vector.parentzone
                self.active_zones.expand_tree(self.active_zone)

        elif self.action == 'select node by node':
            currray: WolfArray
            curarray = self.active_array
            curarray.mngselection.add_node_to_selection(x, y)
            curarray.mngselection.update_nb_nodes_sections()
            # self.Paint()
        elif 'select by tmp vector' in self.action or 'select by vector' in self.action:
            self.active_vector.add_vertex(wolfvertex(x, y))
        elif self.action == 'capture vertices':
            self.active_vector.add_vertex(wolfvertex(x, y))
        elif self.action == 'dynamic parallel':

            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x, y)
                    self.active_vector.myvertices[-1].z = z

            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_zone.parallel_active(self.dynapar_dist)
        elif self.action == 'modify vertices':
            if self.active_vertex is None:
                self.active_vertex = self.active_vector.find_nearest_vert(x, y)
            else:
                self.active_vertex.limit2bounds(self.active_vector.mylimits)
                self.active_vertex = None
        elif self.action == 'insert vertices':
            if self.active_vertex is None:
                self.active_vertex = self.active_vector.insert_nearest_vert(x, y)
            else:
                self.active_vertex = None
        else:
            self.rightdown = (x, y)
        
        # self.Refresh()

    def OnRightUp(self, e):
        pos = e.GetPosition()
        x, y = self.getXY(pos)

        try:
            minx = min(self.rightdown[0], x)
            miny = min(self.rightdown[1], y)
            maxx = max(self.rightdown[0], x)
            maxy = max(self.rightdown[1], y)

            if minx != maxx and maxy != miny:
                self.mybc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'X')
                self.mybc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'Y')
            else:
                self.mybc.query_kdtree((x, y))

            self.mybc.Populate()
            self.Refresh()

        except:
            pass

    def OnButton(self, e: wx.MouseEvent):
        d = e.GetWheelDelta()
        r = e.GetWheelRotation()
        a = e.GetWheelAxis()

        altdown = e.AltDown()
        ctrldown = e.ControlDown()
        shiftdown = e.ShiftDown()

        if self.action == 'dynamic parallel' and shiftdown and not ctrldown:
            self.dynapar_dist *= (1 - .1 * (r / max(d, 1)))
            self.dynapar_dist = max(self.dynapar_dist, .01)

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return
        elif self.action == 'dynamic parallel' and shiftdown and ctrldown:
            dlg = wx.NumberEntryDialog(None, _('What is the desired size [cm] ?'), 'ds', 'ds size',
                                       int(self.dynapar_dist * 100.), 1, 100000.)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self.dynapar_dist = float(dlg.GetValue()) / 100.
            self.dynapar_dist = max(self.dynapar_dist, .01)
            dlg.Destroy()

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return

        self.width = self.width * (1 - .1 * (r / max(d, 1)))
        self.height = self.height * (1 - .1 * (r / max(d, 1)))

        self.setbounds()

    def OnRDClick(self, e):
        self.endactions()

    def OnLDClick(self, e):
        pos = e.GetPosition()
        x, y = self.getXY(pos)
        self.mousex = self.mousedown[0]
        self.mousey = self.mousedown[1]
        self.mousedown = (0., 0.)
        self.oneclick = False
        self.setbounds()

    def OnLDown(self, e):
        if not self.move:
            pos = e.GetPosition()
            x, y = self.getXY(pos)
            self.mousedown = (x, y)
            self.move=True

    def OnActivateTreeElem(self, e): #:dataview.TreeListEvent ):
        curzones: Zones
        curzone: zone
        curvect: vector

        myitem = e.GetItem()
        ctrl = wx.GetKeyState(wx.WXK_CONTROL)

        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        

        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        myobj = self.treelist.GetItemData(myitem)
        self.selobj = myobj
        self.selitem.SetLabel(nameitem)

        if type(myobj) == Zones:
            if ctrl:
                myobj.showstructure(self)
                                
        elif type(myobj) in [WolfArray,WolfArray_Sim2D,WolfArrayMB]:
            if ctrl:
                myobj.myops.SetTitle(_('Operations on array: ')+myobj.idx)
                myobj.myops.Show()
                
            wx.LogMessage(_('Activating array : ' + nameitem))
            self.active_array = myobj
            
        elif type(myobj) == cloud_vertices:
            if ctrl:
                myobj.myprop.show()
        elif type(myobj) == crosssections:
            if ctrl:
                myobj.showstructure(self)
            wx.LogMessage(_('Activating cross sections : ' + nameitem))
            self.active_cs = myobj
        elif type(myobj) == Wolfresults_2D:
            wx.LogMessage(_('Activating Wolf2d results : ' + nameitem))
            self.active_res2d = myobj
            
            if ctrl:
                dlg = wx.MessageDialog(self,_('Would you like to open the 2D model?'),style=wx.YES_NO|wx.NO_DEFAULT)
                ret=dlg.ShowModal()
                dlg.Destroy()
                if ret == wx.ID_NO:
                    return
                
                from .PyGui import Wolf2DModel
                mywolf = Wolf2DModel(os.path.dirname(self.active_res2d.filenamegen))
            
            # myobj.myblocks['block1'].current.mypal.export_image('d:\palh.png')


    def OnMotion(self, e: wx.MouseEvent):

        # Déplacement de la souris sur le canvas OpenGL
        posframe = self.GetPosition()
        pos = e.GetPosition()
        x, y = self.getXY(pos)
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()
        
        if e.LeftIsDown() or e.MiddleIsDown():
            self.mousex -= x - self.mousedown[0]
            self.mousey -= y - self.mousedown[1]

            self.setbounds(False)
            return
        else:
            self.move=False

        if self.action is not None:

            if 'select by tmp vector' in self.action or \
               'select by vector' in self.action or \
               self.action == 'capture vertices' or \
               self.action == 'dynamic parallel':
                
                if self.active_vector.nbvertices > 0:
                    self.active_vector.myvertices[-1].x = x
                    self.active_vector.myvertices[-1].y = y

            if self.action == 'modify vertices' or \
               self.action == 'insert vertices':
                if self.active_vertex is not None:
                    if shiftdown:
                        ox = self.active_vector.myvertices[0].x
                        oy = self.active_vector.myvertices[0].y

                        dirx = self.active_vector.myvertices[-1].x - ox
                        diry = self.active_vector.myvertices[-1].y - oy
                        normdir = np.sqrt(dirx ** 2. + diry ** 2.)

                        dirx /= normdir
                        diry /= normdir

                        vecx = x - ox
                        vecy = y - oy

                        norm = np.sqrt(vecx ** 2. + vecy ** 2.)

                        self.active_vertex.x = ox + np.inner([dirx, diry], [vecx, vecy]) * dirx
                        self.active_vertex.y = oy + np.inner([dirx, diry], [vecx, vecy]) * diry

                    else:
                        self.active_vertex.x = x
                        self.active_vertex.y = y
                        
                    self.active_vertex.limit2bounds(self.active_vector.mylimits)

            if self.action == 'dynamic parallel':
                self.active_zone.parallel_active(self.dynapar_dist)

        self.mytooltip.myparams.clear()

        curgroup = 'Position'
        self.mytooltip.myparams[curgroup] = {}

        curpar = 'Pixel'
        self.mytooltip.myparams[curgroup][curpar] = {}
        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Pixel'
        self.mytooltip.myparams[curgroup][curpar]['value'] = '(' + str(pos[0]) + ' ; ' + str(pos[1]) + ')'
        self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

        curpar = 'x'
        self.mytooltip.myparams[curgroup][curpar] = {}
        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Coordinate X'
        self.mytooltip.myparams[curgroup][curpar]['value'] = '{:3f}'.format(x)
        self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

        curpar = 'y'
        self.mytooltip.myparams[curgroup][curpar] = {}
        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Coordinate Y'
        self.mytooltip.myparams[curgroup][curpar]['value'] = '{:3f}'.format(y)
        self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

        for locarray in self.myres2D:
            curgroup = locarray.idx
            if self.added['wolf2d'][curgroup]['checked']:
                try:
                    val = locarray.get_value(x, y, True)
                    i, j, curbloc = locarray.get_blockij_from_xy(x, y, False)

                    if i != '-':
                        self.mytooltip.myparams[curgroup] = {}
                        curpar = 'Indices'
                        self.mytooltip.myparams[curgroup][curpar] = {}
                        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Indice (i;j;bloc)'
                        self.mytooltip.myparams[curgroup][curpar]['value'] = '(' + str(i + 1) + ';' + str(
                            j + 1) + ';' + str(curbloc) + ')'
                        self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

                        curpar = 'Value'
                        self.mytooltip.myparams[curgroup][curpar] = {}
                        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Value'
                        self.mytooltip.myparams[curgroup][curpar]['value'] = float(val)
                        self.mytooltip.myparams[curgroup][curpar]['type'] = 'Float'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''
                except:
                    pass

        for locarray in self.myarrays:
            curgroup = locarray.idx
            if self.added['arrays'][curgroup]['checked']:

                try:
                    val = locarray.get_value(x, y)

                    if val != -99999.:
                        self.mytooltip.myparams[curgroup] = {}

                        if locarray.wolftype in WOLF_ARRAY_MB:
                            i, j, curbloc = locarray.get_blockij_from_xy(x, y)
                            curpar = 'Indices'
                            self.mytooltip.myparams[curgroup][curpar] = {}
                            self.mytooltip.myparams[curgroup][curpar]['name'] = 'Indice (i;j;bloc)'
                            self.mytooltip.myparams[curgroup][curpar]['value'] = '(' + str(i + 1) + ';' + str(
                                j + 1) + ';' + str(curbloc) + ')'
                            self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
                            self.mytooltip.myparams[curgroup][curpar]['comment'] = ''
                        else:
                            i, j = locarray.get_ij_from_xy(x, y)
                            curpar = 'Indices'
                            self.mytooltip.myparams[curgroup][curpar] = {}
                            self.mytooltip.myparams[curgroup][curpar]['name'] = 'Indice (i;j)'
                            self.mytooltip.myparams[curgroup][curpar]['value'] = '(' + str(i + 1) + ';' + str(
                                j + 1) + ')'
                            self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
                            self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

                        curpar = 'Value'
                        self.mytooltip.myparams[curgroup][curpar] = {}
                        self.mytooltip.myparams[curgroup][curpar]['name'] = 'Value'
                        self.mytooltip.myparams[curgroup][curpar]['value'] = val
                        self.mytooltip.myparams[curgroup][curpar]['type'] = 'Float'
                        self.mytooltip.myparams[curgroup][curpar]['comment'] = ''
                except:
                    pass

        if self.linked:
            for curFrame in self.linkedList:
                if not curFrame is self:
                    for locarray in curFrame.myarrays:
                        curgroup = locarray.idx
                        if curFrame.added['arrays'][curgroup]['checked']:
                            self.mytooltip.myparams[curgroup] = {}

                            try:
                                val = locarray.get_value(x, y)
                                i, j = locarray.get_ij_from_xy(x, y)

                                curpar = 'Indices'
                                self.mytooltip.myparams[curgroup][curpar] = {}
                                self.mytooltip.myparams[curgroup][curpar]['name'] = 'Indice (i;j)'
                                self.mytooltip.myparams[curgroup][curpar]['value'] = '(' + str(i + 1) + ';' + str(
                                    j + 1) + ')'
                                self.mytooltip.myparams[curgroup][curpar]['type'] = 'String'
                                self.mytooltip.myparams[curgroup][curpar]['comment'] = ''

                                curpar = 'Value'
                                self.mytooltip.myparams[curgroup][curpar] = {}
                                self.mytooltip.myparams[curgroup][curpar]['name'] = 'Value'
                                self.mytooltip.myparams[curgroup][curpar]['value'] = val
                                self.mytooltip.myparams[curgroup][curpar]['type'] = 'Float'
                                self.mytooltip.myparams[curgroup][curpar]['comment'] = ''
                            except:
                                pass

        self.mytooltip.PopulateOnePage()
        self.Refresh()

        if e.ControlDown():
            ttsize = self.mytooltip.GetSize()
            self.mytooltip.position(pos + posframe + (ttsize[0] / 2. + 15, 15))
        else:
            width, height = self.GetSize()
            posframe[0] += width
            self.mytooltip.position(posframe)

        self.mytooltip.Show(True)

    def Autoscale(self, update_backfore=True):
        self.findminmax()
        self.width = self.xmax - self.xmin
        self.height = self.ymax - self.ymin

        centerx = self.xmin + self.width / 2.
        centery = self.ymin + self.height / 2.

        iwidth = self.width * self.sx
        iheight = self.height * self.sy

        width, height = self.canvas.GetSize()

        sx = float(width) / float(iwidth)
        sy = float(height) / float(iheight)

        if sx > sy:
            self.xmax = self.xmin + self.width * sx / sy
            self.width = self.xmax - self.xmin
        else:
            self.ymax = self.ymin + self.height * sy / sx
            self.height = self.ymax - self.ymin

        self.mousex = centerx
        self.mousey = centery

        if update_backfore:
            # dessin du background
            for id, (key, curmap) in enumerate(self.added['wms-background'].items()):
                if curmap['checked']:
                    curmap['values'].reload()
            # dessin du foreground
            for id, (key, curmap) in enumerate(self.added['wms-foreground'].items()):
                if curmap['checked']:
                    curmap['values'].reload()

        self.setbounds()

    def endactions(self):
        
        if self.action is not None:
            if 'select by tmp vector' in self.action or 'select by vector' in self.action:
                inside_under = 'inside' in self.action
                self.action = None

                self.active_vector.nbvertices = self.active_vector.nbvertices - 1
                self.active_vector.myvertices.pop(-1)

                if inside_under:
                    self.active_vector.close_force()
                    self.active_array.mngselection.select_insidepoly(self.active_vector)
                else:
                    self.active_array.mngselection.select_underpoly(self.active_vector)

            if self.action == 'capture vertices':
                self.action = None
                self.active_vector.nbvertices = self.active_vector.nbvertices - 1
                self.active_vector.myvertices.pop(-1)
                r = wx.MessageDialog(
                    None,
                    _('End of points capturing') + '\n' +
                    _('Force to close the vector ?'),
                    _('Confirm'),
                    wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION
                ).ShowModal()
                if r == wx.ID_YES:
                    self.active_vector.close_force()

            if self.action == 'modify vertices':
                self.active_vertex = None
                self.action = None

            if self.action == 'insert vertices':
                self.active_vertex = None
                self.action = None

            if self.action == 'dynamic parallel':
                self.active_vector.nbvertices -= 1
                self.active_vector.myvertices.pop(-1)
                self.active_zone.parallel_active(self.dynapar_dist)

                self.active_zones.fill_structure()

                self.active_vertex = None
                self.action = None

            if self.action == 'select active vector' or self.action == 'select active vector2' or self.action == 'select node by node':
                self.action = None

        self.copyfrom = None

        self.mimicme()
                    
    def OnHotKey(self, e: wx.KeyEvent):
        '''
        Gestion des touches clavier
            F2 : mise à jour du résultat pas suivant
            F5 : autoscale
            F7 : refresh
            Z  : zoom avant
            z  : zoom artrière
            Flèches : déplacements latéraux
            P : sélection de profil

            !! ACTIONs !!
            N : sélection noeud par noeud de la matrice courante
            B : sélection par vecteur temporaire de la matrice courante
            V : sélection par vecteur activé de la matrice courante
            RETURN : fin d'action (cf aussi double clicks droit 'OnRDClick')
        '''
        key = e.GetKeyCode()
        ctrldown = e.ControlDown()
        altdown = e.AltDown()

        wx.LogMessage(_('You are pressing key code : ') + str(key))
        if ctrldown:
            wx.LogMessage(_('Ctrl is down'))
        if altdown:
            wx.LogMessage(_('Alt is down'))
        
        if ctrldown or altdown:
            if key == ord('U'):
                # CTRL+U
                # Mise à jour des données par import du fichier gtlf2
                msg = ''
                if self.active_array is None:
                    msg += _('Active array is None\n')

                if msg != '':
                    msg += _('\n')
                    msg += _('Retry !\n')
                    wx.MessageBox(msg)
                    return

                self.set_fn_fnpos_gltf()
                self.update_blender_sculpting()
            elif key == ord('C') and e.ControlDown():
                # CTRL+C 
                # CTRL+ALT+C ou Alt Gr + C
                
                if self.active_array is None:
                    if e.AltDown():
                        dlg = wx.MessageDialog(self,
                                            _('The active array is None - Please active an array from which to copy the selection !'),
                                            style=wx.OK)                   
                    else:
                        dlg = wx.MessageDialog(self,
                                            _('The active array is None - Please active an array from which to copy the values !'),
                                            style=wx.OK)
                    dlg.ShowModal()
                    return

                wx.LogMessage(_('Copy selection'))
                self.copyfrom = self.active_array
                self.mimicme_copyfrom()  # force le recopiage de copyfrom dans les autres matrices liées

            elif key == ord('V') and e.ControlDown():
                # CTRL+V
                # CTRL+ALT+V ou Alt Gr + V
                
                if self.active_array is None:
                    if e.AltDown():
                        # CTRL+ALT+V
                        wx.LogWarning(_('The active array is None - Please active an array into which to paste the selection !'))
                    else:
                        wx.LogWarning(_('The active array is None - Please active an array into which to paste the values !'))
                        
                    return
                        
                fromarray = self.copyfrom
                if fromarray is None:
                    if self.linked:
                        if not self.linkedList is None:
                            for curFrame in self.linkedList:
                                if curFrame.copyfrom is not None:
                                    fromarray = curFrame.copyfrom
                                    break
                
                if fromarray is None:
                    wx.LogWarning(_('No selection to be pasted !'))
                    return

                cursel = fromarray.mngselection.myselection

                if e.AltDown():
                    wx.LogMessage(_('Paste selection position'))
                    
                    if cursel == 'all':
                        self.active_array.mngselection.OnAllSelect(0)
                    elif len(cursel) > 0:
                        self.active_array.mngselection.myselection = cursel.copy()
                        self.active_array.mngselection.update_nb_nodes_sections()                                
                else:
                    wx.LogMessage(_('Copy selection values'))
                    if cursel == 'all':
                        self.active_array.paste_all(fromarray)
                    elif len(cursel) > 0:
                        z = fromarray.mngselection.get_values_sel()
                        self.active_array.set_values_sel(cursel, z)    
        else:
            if key == 13 or key==370:
                # 13 = RETURN classic keyboard
                # 370 = RETURN NUMPAD
                self.endactions()

            elif key == wx.WXK_F2:
                # Dessin des matrices
                try:
                    for curarray in self.added['arrays']:
                        if self.added['arrays'][curarray]['checked']:
                            if type(self.added['arrays'][curarray]['values']) is wolfres2DGPU:
                                self.added['arrays'][curarray]['values'].onnext(1)

                    self.Paint()
                except:
                    pass
            elif key == wx.WXK_F5:
                # Autoscale
                self.Autoscale()
            elif key == wx.WXK_F7:
                self.update()
            elif key == wx.WXK_F12:
                self.active_array.myops.SetTitle(_('Operations on array: ')+self.active_array.idx)
                self.active_array.myops.Show()
                self.active_array.myops.array_ops.SetSelection(1)
                self.active_array.myops.Center()
            elif key == ord('N'):  # N
                if self.active_array is not None:
                    self.active_array.myops.select_nod_by_node()
            elif key == ord('V'):  # V
                if self.active_array is not None:
                    self.active_array.myops.select_vector_inside_manager()
            elif key == ord('B'):  # B
                if self.active_array is not None:
                    self.active_array.myops.select_vector_tmp()
            elif key == ord('P'):  # P
                self.action = 'Select nearest profile'
            elif key == ord('Z'):  # Z
                self.width = self.width / 1.1
                self.height = self.height / 1.1
                self.setbounds()
            elif key == ord('z'):  # z
                self.width = self.width * 1.1
                self.height = self.height * 1.1
                self.setbounds()
            elif key == wx.WXK_UP:
                self.mousey = self.mousey + self.height / 10.
                self.setbounds()
            elif key == wx.WXK_DOWN:
                self.mousey = self.mousey - self.height / 10.
                self.setbounds()
            elif key == wx.WXK_LEFT:
                self.mousex = self.mousex - self.width / 10.
                self.setbounds()
            elif key == wx.WXK_RIGHT:
                self.mousex = self.mousex + self.width / 10.
                self.setbounds()                

    def paste_values(self,fromarray:WolfArray):
        wx.LogMessage(_('Copy selection values'))
        cursel = fromarray.mngselection.myselection
        if cursel == 'all':
            self.active_array.paste_all(fromarray)
        elif len(cursel) > 0:
            z = fromarray.mngselection.get_values_sel()
            self.active_array.set_values_sel(cursel, z)            

    def paste_selxy(self,fromarray:WolfArray):
        wx.LogMessage(_('Paste selection position'))
        cursel = fromarray.mngselection.myselection
        if cursel == 'all':
            self.active_array.mngselection.OnAllSelect(0)
        elif len(cursel) > 0:
            self.active_array.mngselection.myselection = cursel.copy()
            self.active_array.mngselection.update_nb_nodes_sections()                                

    def OntreeRight(self, e: wx.MouseEvent):
        if self.selobj is not None:
            self.treelist.PopupMenu(self.popupmenu)

    def update(self):
        # dessin du background
        for id, (key, curmap) in enumerate(self.added['wms-background'].items()):
            if curmap['checked']:
                curmap['values'].reload()
        # dessin du foreground
        for id, (key, curmap) in enumerate(self.added['wms-foreground'].items()):
            if curmap['checked']:
                curmap['values'].reload()

        if self.locminmax.IsChecked() or self.update_absolute_minmax:
            for id, (key, curmap) in enumerate(self.added['arrays'].items()):
                if curmap['checked']:
                    curarray: WolfArray
                    curarray = curmap['values']
                    if self.update_absolute_minmax:
                        curarray.updatepalette()
                        self.update_absolute_minmax = False
                    else:
                        curarray.updatepalette(onzoom=[self.xmin, self.xmax, self.ymin, self.ymax])
                    curarray.delete_lists()
        self.Refresh()

    def Paint(self):

        if self.currently_readresults:
            return

        width, height = self.canvas.GetSize()

        # C'est bien ici que la zone de dessin utile est calculée sur base du centre et de la zone en coordonnées réelles
        # Les commandes OpenGL sont donc traitées en coordonnées réelles puisque la commande glOrtho définit le cadre visible
        self.xmin = self.mousex - self.width / 2.
        self.ymin = self.mousey - self.height / 2.
        self.xmax = self.mousex + self.width / 2.
        self.ymax = self.mousey + self.height / 2.

        if self.canvas.SetCurrent(self.context):

            glClearColor(1., 1., 1., 0)
            glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT)
            glViewport(0, 0, int(width), int(height))

            glMatrixMode(GL_PROJECTION)
            glLoadIdentity()
            glOrtho(self.xmin, self.xmax, self.ymin, self.ymax, -99999, 99999)

            glMatrixMode(GL_MODELVIEW)
            glLoadIdentity()

            # dessin du background
            for id, (key, curmap) in enumerate(self.added['wms-background'].items()):
                if curmap['checked']:
                    curmap['values'].paint()

            # Dessin des matrices
            try:
                for curarray in self.added['arrays']:
                    if self.added['arrays'][curarray]['checked']:
                        locarray = self.added['arrays'][curarray]['values']
                        if not locarray.plotting:
                            locarray.plotting = True
                            self.added['arrays'][curarray]['values'].plot(self.sx, self.sy, self.xmin, self.ymin,
                                                                          self.xmax, self.ymax)
                            locarray.plotting = False
            except:
                pass

            # Dessin des résultats 2D
            try:
                for curarray in self.added['wolf2d']:
                    if self.added['wolf2d'][curarray]['checked']:
                        locarray = self.added['wolf2d'][curarray]['values']
                        if not locarray.plotting:
                            locarray.plotting = True
                            self.added['wolf2d'][curarray]['values'].plot(self.sx, self.sy, self.xmin, self.ymin,
                                                                          self.xmax, self.ymax)
                            locarray.plotting = False
            except:
                pass

            # Dessin des vecteurs
            try:
                for curvect in self.added['vectors']:
                    if self.added['vectors'][curvect]['checked']:
                        self.added['vectors'][curvect]['values'].plot()
            except:
                pass

            # Dessin des nuages
            try:
                for curcloud in self.added['clouds']:
                    if self.added['clouds'][curcloud]['checked']:
                        self.added['clouds'][curcloud]['values'].plot()
            except:
                pass

            # Dessin du reste
            try:
                for curobj in self.added['others']:
                    if self.added['others'][curobj]['checked']:
                        curobj = self.added['others'][curobj]['values']
                        if type(curobj) is SPWDCENNGaugingStations or type(curobj) is SPWMIGaugingStations:
                            curobj.plot((self.xmax - self.xmin) / 100.)
                        elif type(curobj) is genericImagetexture:
                            curobj.paint()

            except:
                pass

            # Dessin du Front
            for id, (key, curmap) in enumerate(self.added['wms-foreground'].items()):
                if curmap['checked']:
                    curmap['values'].paint()

            # Gestion des BC (si actif)
            try:
                if self.mybc != None:
                    self.mybc.plot()
            except:
                pass

            glFlush()
            self.canvas.SwapBuffers()
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def OnPaint(self, e):

        self.Paint()
        if e != None:
            e.Skip()

    def findminmax(self, force=False):
        xmin = 1.e30
        ymin = 1.e30
        xmax = -1.e30
        ymax = -1.e30

        k = 0
        for locarray in self.myarrays:
            if self.added['arrays'][locarray.idx]['checked'] or force:
                xmin = min(locarray.origx + locarray.translx, xmin)
                xmax = max(locarray.origx + locarray.translx + float(locarray.nbx) * locarray.dx, xmax)
                ymin = min(locarray.origy + locarray.transly, ymin)
                ymax = max(locarray.origy + locarray.transly + float(locarray.nby) * locarray.dy, ymax)
                k += 1

        for locvector in self.myvectors:
            if self.added['vectors'][locvector.idx]['checked'] or force:
                if locvector.idx != 'grid':
                    locvector.find_minmax()
                    xmin = min(locvector.minx, xmin)
                    xmax = max(locvector.maxx, xmax)
                    ymin = min(locvector.miny, ymin)
                    ymax = max(locvector.maxy, ymax)
                    k += 1

        for loccloud in self.myclouds:
            if self.added['clouds'][loccloud.idx]['checked'] or force:
                loccloud.find_minmax(force)
                xmin = min(loccloud.xbounds[0], xmin)
                xmax = max(loccloud.xbounds[1], xmax)
                ymin = min(loccloud.ybounds[0], ymin)
                ymax = max(loccloud.ybounds[1], ymax)
                k += 1

        for locres2d in self.myres2D:
            if self.added['wolf2d'][locres2d.idx]['checked'] or force:
                xmin = min(locres2d.origx + locres2d.translx, xmin)
                xmax = max(locres2d.origx + locres2d.translx + float(locres2d.nbx) * locres2d.dx, xmax)
                ymin = min(locres2d.origy + locres2d.transly, ymin)
                ymax = max(locres2d.origy + locres2d.transly + float(locres2d.nby) * locres2d.dy, ymax)
                k += 1

        for locothers in self.myothers:
            if type(locothers) is genericImagetexture:
                xmin = locothers.xmin
                xmax = locothers.xmax
                ymin = locothers.ymin
                ymax = locothers.ymax
                k += 1

        if k > 0:
            self.xmin = xmin
            self.xmax = xmax
            self.ymin = ymin
            self.ymax = ymax

    def resizeFrame(self, w, h):
        self.SetClientSize(w, h)

    def mimicme(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.resizeFrame(width, height)
                        curFrame.mousex = self.mousex
                        curFrame.mousey = self.mousey
                        curFrame.sx = self.sx
                        curFrame.sy = self.sy
                        curFrame.width = self.width
                        curFrame.height = self.height
                        curFrame.setbounds()

                        if curFrame.link_shareopsvect:
                            curFrame.Active_vector(self.active_vector)
                            curFrame.active_array.myops.Active_vector(self.active_vector, False)
                            curFrame.action = self.action

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def mimicme_copyfrom(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.copyfrom = self.copyfrom

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def Active_vector(self, vect):
        self.active_vector = vect
        
        if vect is not None:
            wx.LogMessage(_('Activating vector : ' + vect.myname))            
            if vect.parentzone is not None:
                self.Active_zone(vect.parentzone)

        self.mimicme()

    def Active_zone(self, zone: zone):
        self.active_zone = zone
        self.active_zones = zone.parent
        wx.LogMessage(_('Activating zone : ' + zone.myname))            



# importing all necessary modules
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import pandas as pd
import re
import nltk
from nltk.stem.snowball import SnowballStemmer
st = SnowballStemmer('english')
from nltk.stem import PorterStemmer
from nltk.stem.wordnet import WordNetLemmatizer
from transformers import AutoModelForSequenceClassification, AutoTokenizer, pipeline
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from datetime import date


class Text_Labling:
    
    def __init__(self,path,column,clean_col,stem):
        print('Loading Data...')
        print('Done..\n\n')
        self.stop_words= nltk.corpus.stopwords.words("english")
        self.data= pd.read_csv(path)
        self.column=column
        self.clean_col=clean_col
        self.stem=stem
        
    def cleanData(self):
        #Getting English Text Only
        print('Cleaning Data...')
        print('Done..\n\n')
        self.data['Language'] = self.data[self.column].apply(lambda S: self.isEnglish(S))
        self.data=self.data[self.data['Language']=='en']
        
        # change to lower and remove spaces on either side
        self.data[ self.clean_col] = self.data[self.column].apply(lambda x: x.lower().strip())
    
        # remove extra spaces in between
        self.data[ self.clean_col] = self.data[ self.clean_col].apply(lambda x: re.sub(' +', ' ', x))
    
        # remove punctuation
        self.data[ self.clean_col] = self.data[ self.clean_col].apply(lambda x: re.sub('[^a-zA-Z]', ' ', x))
    
        # remove stopwords and get the stem
        self.data[ self.clean_col] = self.data[ self.clean_col].apply(lambda x: ' '.join(st.stem(text) for text in x.split() if text not in self.stop_words))
        
        if self.stem == 'Stem':
             stemmer = PorterStemmer() 
             self.data[ self.clean_col] = [stemmer.stem(y) for y in self.data[ self.clean_col]]
        elif self.stem == 'Lem':
             lem = WordNetLemmatizer()
             self.data[ self.clean_col] = [lem.lemmatize(y) for y in self.data[ self.clean_col]]
        # elif self.stem == 'Spacy':
        #      self.data[ self.clean_col] = spac(' '.join(self.data[self.clean_col]))
        #      self.data[ self.clean_col] = [y.lemma_ for y in self.data[self.clean_col]]
        elif self.stem == 'None':
             self.data[ self.clean_col]==self.data[ self.clean_col]
        
        
      
        return self.data
    
    def isEnglish(self,s):
        s = str(s)
        try:
            s.encode(encoding='utf-8').decode('ascii')
        except UnicodeDecodeError:
            return "ar"
        else:
            return "en"
    
    
    def DoLabels(self,labels):
        print('Loading Bart Model..')
        # calling function
        bart_model = AutoModelForSequenceClassification.from_pretrained('navteca/bart-large-mnli')
        bart_tokenizer = AutoTokenizer.from_pretrained('navteca/bart-large-mnli')
        print('Done..\n\n')
        
        # Get predictions
        self.data.reset_index(inplace=False)
        
        
        labels=labels
        nlp = pipeline('zero-shot-classification', model=bart_model, tokenizer=bart_tokenizer)

        print('Looking for nulls..')
        
        self.data[self.clean_col].replace('', np.nan, inplace=True)
        print('Total Nulls: ',self.data[self.clean_col].isnull().sum())
        print('\nDropping nulls..')
        self.data.dropna(subset=[self.clean_col], inplace=True)
        print('Nulls Dropped..')
        print('\nLabling Started...')
        
        
        lab=[]
        word=[]
        score=[]
        c=0
        for i in self.data[self.clean_col]:
            print(i,c)
            c+=1
            sequence= i
            x=nlp(sequence, labels)
            first_value = list(x.items())[0][1]
            word.append(first_value)
            seconed_value = (list(x.items())[1][1][0])
            lab.append(seconed_value)
            Third_value = list(x.items())[2][1][0]
            score.append(Third_value)
        print('\n\nDone Labling:\n=======================  \n\n\n  ')
        se = pd.Series(lab)
        so = pd.Series(score)

        self.data['Words']=word
        self.data['Labels']=se.values
        self.data['Score']=so.values     
        #Results
        print('Dataset head:\n=======================  \n\n\n  ',self.data.head())
        print('Dataset info:\n=======================  \n\n\n  ',self.data.info())
        self.data.to_csv('datalabeld.csv',index=False)
    
    
    def doVis(self,num):
              #plot the Word Pie chart image 
              #self.data['Labels'].value_counts().plot(kind='pie')
              
              
              # wordcloud
              comment_words = ''
              stop_words= nltk.corpus.stopwords.words("english")
              stopwords = stop_words
    
              # iterate through the csv file
              for val in self.data['Words']:
                   
                  # typecaste each val to string
                  val = str(val)
               
                  # split the value
                  tokens = val.split()
                   
                  # Converts each token into lowercase
                  for i in range(len(tokens)):
                      tokens[i] = tokens[i].lower()
                   
                  comment_words += " ".join(tokens)+" "
               
              wordcloud = WordCloud(width = 800, height = 800,
                              background_color ='white',
                              stopwords = stopwords,
                              min_font_size = 10).generate(comment_words)
               
              # plot the WordCloud image                      
              plt.figure(figsize = (10, 10), facecolor = None)
              plt.imshow(wordcloud)
              plt.axis("off")
              plt.tight_layout(pad = 0)
              plt.savefig("WordCloud.png")
              plt.show()

              # plot pie chart
              self.data['Labels'].value_counts().plot(kind='pie',colors = ['#1B204C','#1BB790','#F0BD5B']
                                                      ,shadow=True,explode=(0.1, 0.1, 0.1, 0.1), autopct='%1.1f%%')
              plt.title("Lables Distribution")
              plt.savefig("Lables Distribution.png")
              plt.show()
         
              
              #plot the Top frequent image 
              top_N = num
              txt = self.data['Words'].str.replace(r'\|', ' ').str.cat(sep=' ')
              words = nltk.tokenize.word_tokenize(txt)            
              words_except_stop_dist = nltk.FreqDist(w for w in words if w not in self.stop_words) 
              rslt = pd.DataFrame(words_except_stop_dist.most_common(top_N),
                                columns=['Word', 'Frequency']).set_index('Word')
              rslt.plot.bar(rot=90,color =  ['#F0BD5B'])
              plt.rcParams.update({'axes.facecolor':'#1BB790'})
              plt.title("Frequency of Words")
              plt.xlabel("Frequency")
              plt.ylabel("Words")
              plt.savefig("Top frequent.png")
              plt.show()

              
              # create an Object
              ppt = Presentation()
              first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
               
              # title (included date)
              title = "Bar chart & Wordcloud are added - " + str(date.today())
               
               
              # set the title on first slide
              first_slide.shapes[0].text_frame.paragraphs[0].text = title
               
              # slide 2 - set the image
              img = 'WordCloud.png'
              second_slide = ppt.slide_layouts[1]
              slide2 = ppt.slides.add_slide(second_slide)


              img2 = 'Top frequent.png'
              Third_slide = ppt.slide_layouts[1]
              slide3 = ppt.slides.add_slide(Third_slide)
              
              img3 = 'Lables Distribution.png'
              Fourth_slide = ppt.slide_layouts[1]
              slide4 = ppt.slides.add_slide(Fourth_slide)
              # play with the image attributes if you are not OK with the height and width
              pic = slide2.shapes.add_picture(img, left= Inches(2),top = Inches(1),height = Inches(5))
              pic2 = slide3.shapes.add_picture(img2, left= Inches(2),top = Inches(1),height = Inches(5))
              pic3 = slide4.shapes.add_picture(img3, left= Inches(2),top = Inches(1),height = Inches(5))

              # save the powerpoint presentation
              ppt.save('Bar chart & Wordcloud.pptx')
              print('All Done..')

              

             
                                    
      


    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    


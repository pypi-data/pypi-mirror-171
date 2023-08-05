#Statistics and graphs libraries
import os
import numpy as np
import pandas as pd #datafram functions
import matplotlib.pyplot as plt
from scipy.stats import ttest_ind #ttest functions
from scipy.stats import ttest_rel
from statsmodels.formula.api import ols
from statsmodels.sandbox.stats.multicomp import MultiComparison
from statsmodels.stats.anova import anova_lm
import matplotlib.pyplot as plt
from matplotlib.ticker import FormatStrFormatter
from openpyxl.workbook import Workbook

#Powerpoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import win32com.client


#2: bar graph with unparied t-test for between groups within 1 paramater.
#2.1: 'flt' is to filter out the entire dataframe from a given string
def bar_unttest(df,para,grp,ymax,flth='',flt='',y_min='None',y_label=''):
    #Filter dataframe if flt !=''
    if flth != '':
        if isinstance(flt,list):
            for e in flt:
                df = df[df[flth].isin(flt)]
        else:
            df = df[df[flth]==flt]
   
    # Statistics: mean,sem,count
    msc = df.groupby([grp])[para].agg(['mean','sem','count'])
    msc.reset_index(inplace=True) #converts any columns in index as columns
    pd.DataFrame(msc)
    with pd.ExcelWriter(os.getcwd()+"\msc.xlsx",mode='a',engine='openpyxl',if_sheet_exists='replace') as writer:
        msc.to_excel(writer, sheet_name=para+" barunttest", index=False)
    
    #Statistics: Two-sample unpaired T-TEST
    u1 = msc[grp].unique() #unique values in group 
    cat1 = df[df[grp]==u1[0]]
    cat2 = df[df[grp]==u1[1]]
    results = ttest_ind(cat1[para], cat2[para])
    pval = results[1]
   
    #Graph: numerical calculations
    y1,yerr1 = msc['mean'].to_numpy(),msc['sem'].to_numpy() 
    x1=np.arange(len(y1)) #count the number of values in y that are not zero which will form the number of x-positions
    
    #Graph: Create
    fig, ax = plt.subplots(1,1, figsize = (3,4))
    ax.bar(x1, y1, color="#00004d",width=0.5)
    ax.errorbar(x1, y1, yerr=yerr1,fmt=' ', capsize=(3),zorder=0, ecolor='k')
    ax.spines['right'].set_visible(False), ax.spines['top'].set_visible(False)
    ax.set_xticks(x1)
    ax.set_xticklabels(u1)
    if y_label != '':
        ax.set_ylabel(ylabel)
    else:
        ax.set_ylabel(para.replace('_',' ')+'(g)', labelpad=5) #labels
    if y_min != 'None':
        y_min = y_min
    ax.set_ylim(top=ymax,bottom=y_min)
    ax.set_xlim(-0.5, 1.5)
    ax.set_xlabel(' ')
    ax.locator_params(axis='y', nbins=7) #set the maximum number of ticks on the y-axis
        
    # #save and export graph prior to significance line
    fig.savefig(para+"_JR.png",dpi=300,bbox_inches='tight')
   
    #Set signficance astricks according to pval
    if pval <= 0.0001:
      text = '****'
    elif pval <= 0.001:
      text='***'
    elif pval <=0.01:
       text='**'
    elif pval <= 0.05:
        text='*'
    else:
        text=''
        return
    
    #Create significance annotation
    x = (x1[0]+x1[1])/2
    y = max(y1+yerr1) #ensure error bar sits on top of yerr. This values has to be a single value, hence max
    props = {'connectionstyle':"bar,fraction=0.2",'arrowstyle':'-','shrinkA':10,'shrinkB':10,'linewidth':2} #fraction is the distance of the connecting line from point A and B 
    ax.annotate(text,xy=(x,y*1.16),ha='center',fontsize=15) #text annotation
    ax.annotate('',xy=(x1[0],y),xytext=(x1[1],y),arrowprops=props) #annotates line
    
    #Figure size and export as .png
    fig.savefig(para+"_JR.png",dpi=300,bbox_inches='tight')

#3: Line graph with two-way anova for groups containing 2 subgroups
#Note: 'order' refers to the order of the x axis
def plot_2y(df,grp,xv,yv,order='',ymax=0,x_label='',y_min='None',y_label=''):
    #Statistics: Mean, sem, count
    msc = df.groupby([grp,xv])[yv].agg(['mean','sem','count'])
    msc.reset_index(inplace=True) #converts any columns in index as columns
    pd.DataFrame(msc)
    with pd.ExcelWriter(os.getcwd()+"\msc.xlsx",mode='a',engine='openpyxl',if_sheet_exists='replace') as writer:
        msc.to_excel(writer, sheet_name=yv+" ploy2y", index=False)
    
    #Statistics: TWO-WAY ANOVA AND MULTICOMP
    df['comb'] = df[xv].map(str) + "+" + df[grp].map(str) #add comb column to orginal df
    mod = ols(yv+'~'+grp+'+'+xv+'+'+grp+'*'+xv, data = df).fit()
    aov = anova_lm(mod, type=2) #mod needs to be the same text as mod (i.e. mod1,mod2)
    comparison=MultiComparison(df[yv], df['comb'])
    tdf = pd.read_html(comparison.tukeyhsd().summary().as_html())[0] #tukey's test 
    headings = {'group1':[],'group2':[],'meandiff':[],'p-adj':[],'lower':[],'upper':[],'reject':''}
    data = pd.DataFrame(headings) #Create a new df summary of multicomp results filtered
    for i in order:
        comp = tdf[tdf['group1'].str.startswith(i) & tdf['group2'].str.startswith(i)]
        data = data.append(comp)
 
    #GRAPH
    l1,l2 = df[grp].unique(),df[xv].unique() #labe1 for x-axis and for legends respectfully
    fig = plt.figure()
    fig, ax = plt.subplots(1,1, figsize=(5.5,4)) #figsize alters font size in figure
    y1f,y2f = msc[(msc[grp]==l1[0])],msc[(msc[grp]==l1[1])] #numbers refer to index to find for legend
    y1,y2 = y1f['mean'].to_numpy(),y2f['mean'].to_numpy() #multiassignment
    x1=np.arange(len(y1)) #count the number of ys in y1 that are not zero which will form the number of x-positions
    zlst = [0]*len(y1) #Create a list of zeros according to the number of elements in y1
    yerr1,yerr2 = [zlst,y1f['sem']],[zlst,y2f['sem']] #combined zlst and sem, to create upper errorbars only
    p1 = ax.plot(x1, y1, color='#92414e')
    p2 = ax.plot(x1, y2, color='#200e11') 
    ax.errorbar(x1, y1, yerr=yerr1,capsize=(2),ecolor='k',color='#92414e',fmt='s',ms='4',zorder=2)
    ax.errorbar(x1, y2, yerr=yerr2,capsize=(2),ecolor='k',color='#200e11',fmt='s',ms='4',zorder=2)
    ax.legend([p1, p2], labels=l1,
            loc='upper right',bbox_to_anchor=(1,1),frameon=False)
    ax.spines['right'].set_visible(False)
    ax.spines['top'].set_visible(False)
    ax.set_xticks(x1) #you need this to properly position your xtick labels
    ax.set_xticklabels(l2,ha='center')
    
    if y_label != '':
        ax.set_ylabel(y_label)
    else:
        ax.set_ylabel(yv.replace('_',' ')+'(g)', labelpad=5) #labels
    if x_label != '': #allows adjusting x-axis label if needed
        x_label = x_label
        
    if y_min != 'None':
        y_min = y_min

    ax.set_xlabel(x_label)
    ax.set_ylim(top=ymax,bottom=y_min)
        
    #save and export graph prior to significance line
    fig.savefig(grp+'_'+yv+'_JR.png',dpi=200)
   
    #Significance
    pval = data['p-adj'].tolist()
    print(pval)
    for e in pval:
        round(e,4)
    yn1,yn2= np.array(y1f['sem'])+y1,np.array(y2f['sem'])+y2 #yerr is 2D array, get yerr from original y1f
    ymerge = np.column_stack((y1,y2)) #y values from both grps where yerr and y were added together
    
    for p in pval:
        if p < 0.0001:
            text = '****'
            ax.text(s=text,fontsize=15,ha='center',va='center',x=pval.index(p),y=np.max(ymerge[pval.index(p)])*1.15)
        elif p < 0.001:
            text='***'
            ax.text(s=text,fontsize=15,ha='center',va='center',x=pval.index(p),y=np.max(ymerge[pval.index(p)])*1.15)
        elif p<0.01:
            text='**'
            ax.text(s=text,fontsize=15,ha='center',va='center',x=pval.index(p),y=np.max(ymerge[pval.index(p)])*1.15)
        elif p < 0.05:
            text='*'
            ax.text(s=text,fontsize=15,ha='center',va='center',x=pval.index(p),y=np.max(ymerge[pval.index(p)])*1.15)
        else:
            text=''
                
    #Figure size and export as .png
    fig.savefig(grp+'_'+yv+'_JR.png',dpi=200,bbox_inches='tight')

#4: cluster bar graph but only an paired t-test for within groups within 1 paramater  
def cbar_pttest(df, para, grp, legend, ymax,y_label=''):
    #Statistics: mean, sem, count (msc)
    msc = df.groupby([grp,legend])[para].agg(['mean','sem','count']) #for each combination between grp and legend
    msc.reset_index(inplace=True) #converts any columns in index as columns
    pd.DataFrame(msc)
    with pd.ExcelWriter(os.getcwd()+"\msc.xlsx",mode='a',engine='openpyxl',if_sheet_exists='replace') as writer:
        msc.to_excel(writer, sheet_name=para+" cbarpttest", index=False)

    #Statistics: t.test two tailed unpaired 
    u1,u2 = msc[grp].unique(),msc[legend].unique()#retrieve unique subgroups name from grp (u1) and legend (u2)
    df1,df2 = df[df[grp]==u1[0]],df[df[grp]==u1[1]] #seperate main df to individual df according to u1
    df1,df2,df3,df4 = df1[df1[legend] == u2[0]],df1[df1[legend] == u2[1]],df2[df2[legend] == u2[0]], df2[df2[legend] == u2[1]] #each seperate df seperated according to u2 
    ttest1, ttest2 = ttest_rel(df1[para], df2[para]), ttest_rel(df3[para],df4[para]) #ttest function for grp and legend
    pval1,pval2 = ttest1[1],ttest2[1] #extract pval for significance annotations for graph
    print(para, pval1,pval2)

    #Bar Graph: Numerical Information for graphs
    y1,y2 = msc[msc[legend]==u2[0]],msc[msc[legend]==u2[1]]
    y1,yerr1,y2,yerr2 = y1['mean'].to_numpy(),y1['sem'].to_numpy(),y2['mean'].to_numpy(),y2['sem'].to_numpy()
    x1=np.arange(len(y1)) #x1 and y1 positions need to be equal

    #Bar Graph: Create
    fig = plt.figure()
    fig, ax = plt.subplots(1,1)
    bar_width = 0.3 #seperate two bar graphs
    b1 = ax.bar(x1, y1,width=bar_width, color='#92414e') #bar graph 1
    b2 = ax.bar(x1+bar_width, y2,width=bar_width,color='#200e11') #bar graph 2
    ax.errorbar(x1, y1, yerr=yerr1,fmt=' ', capsize=(3),zorder=0, ecolor='k')
    ax.errorbar(x1+bar_width, y2, yerr=yerr2,fmt=' ', capsize=(3),zorder=0, ecolor='k')
    ax.legend([b1, b2], labels=u2,
            loc='upper right',bbox_to_anchor=(1,1),frameon=False)
    ax.spines['right'].set_visible(False)
    ax.spines['top'].set_visible(False)
    ax.set_xticks(x1+bar_width/2) #you need this to properly position your xtick labels
    ax.set_xticklabels(u1)
    if y_label != '':
        ax.set_ylabel(y_label)
    else:
        ax.set_ylabel(para.replace('_',' ')+'(g)', labelpad=5) #labels
    ax.set_ylim(top=ymax) #y-label but remove '_' and ymax

    #Significance lines for pval 1
    ymerge = np.array([[i, j] for i, j in zip(yerr1+y1, yerr2+y2)]).ravel() #determine height of sig bar by combining b1 and b2 yerr and y values
    if pval1 <= 0.0001:
        text = '****'
        x = (x1[0]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[0],ymerge[1])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[0],y),xytext=(x1[0]+bar_width,y),arrowprops=props) #annotates line
    elif pval1 <= 0.001:
        text='***'
        x = (x1[0]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[0],ymerge[1])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x,y*1.1),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[0],y),xytext=(x1[0]+bar_width,y),arrowprops=props) #annotates line
    elif pval1 <=0.01:
        text='**'
        x = (x1[0]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[0],ymerge[1])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[0],y),xytext=(x1[0]+bar_width,y),arrowprops=props) #annotates line
    elif pval1 <= 0.05:
        text='*'
        x = (x1[0]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[0],ymerge[1])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[0],y),xytext=(x1[0]+bar_width,y),arrowprops=props) #annotates line
    else:
        text=''

    #Significance line for pval2
    if pval2 <= 0.0001:
        text = '****'
        x = (x1[1]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[2],ymerge[3])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x1[1]+0.15,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[1],y),xytext=(x1[1]+bar_width,y),arrowprops=props) #annotates line
    elif pval2 <= 0.001:
        text='***'
        x = (x1[1]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[2],ymerge[3])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x1[1]+0.15,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[1],y),xytext=(x1[1]+bar_width,y),arrowprops=props) #annotates line
    elif pval2 <=0.01:
        text='**'
        x = (x1[1]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[2],ymerge[3])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x1[1]+0.15,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[1],y),xytext=(x1[1]+bar_width,y),arrowprops=props) #annotates line
    elif pval2 <= 0.05:
        text='*'
        x = (x1[1]+bar_width)/2 #in between two bars
        y = 1.01*max(ymerge[2],ymerge[3])
        props = {'connectionstyle':'bar','arrowstyle':'-','shrinkA':5,'shrinkB':5,'linewidth':2}
        ax.annotate(text,xy=(x1[1]+0.15,y*1.08),zorder=10,ha='center',fontsize=15) #text annotation
        ax.annotate('',xy=(x1[1],y),xytext=(x1[1]+bar_width,y),arrowprops=props) #annotates line
    else:
        text=''
    
    #Figure size and export as .png
    fig.set_size_inches(3, 4)
    fig.savefig(para+'_'+grp+'_JR.png',dpi=200, bbox_inches='tight')

#5: Create a figure through ppt
def four_figs(path):
    #Create a blank slide
    global ppt, Presentation
    ppt = Presentation() #Create presentation
    ppt.slide_width,ppt.slide_height = Inches(6), Inches(4.7) #slides dimension
    slide = ppt.slides.add_slide(ppt.slide_layouts[6]) #create a blank slide
    left = top = Inches(0.2) #margins, gap from the top and left where image will be placed

    #Powerpoint slide: import images from folder into slide
    image_files = os.listdir(path) #List files in current directory
    sorted_by_mtime= sorted(image_files, key=lambda t: os.stat(t).st_mtime) #sort in ascending order
    image_files = sorted_by_mtime #replace old image_files with new sorted folder

    n=1 #to iterate through each file in the folder
    for file in image_files: #For loop inputs images into a 3x2 grid (n=6)
        if file.endswith('.png'):
            if n==1 and left<Inches(5.1) and top==Inches(0.2):
                n,top=n+1,top #Position of n image
                img_path = path + "\\" +file #Location of image in folder
                slide.shapes.add_picture(img_path,left,top,height=Inches(1.96)) #Add n image to ppt
            elif n>1 and left<Inches(2) and top==Inches(0.2):
                n,left,top = n+1, left + Inches(2.8),top
                img_path = path + "\\" +file
                slide.shapes.add_picture(img_path,left,top,height=Inches(1.96))
            elif n>1 and left>Inches(2) and top<Inches(2):
                n,left,top=n+1,Inches(0.2),top+Inches(1.96)
                img_path = path + "\\" +file
                slide.shapes.add_picture(img_path,left,top,height=Inches(1.96))
            elif n>1 and left<Inches(2) and top>Inches(1.96):
                n,left,top=n+1,left + Inches(2.8),top
                img_path = path + "\\" +file
                slide.shapes.add_picture(img_path,left,top,height=Inches(1.96))

    #Powerpoint slide: add text to slide per figure
    lb = ['A','B','C','D'] #Each fig has a label
    n,left,top=1,Inches(0.1),Inches(0.1) #Restart these paramaters for labelling

    for letter in lb: #For loop inputs labels next to each figure (n=6)
        if n==1 and left<Inches(5.1) and top==Inches(0.1):
            n,top=n+1,top #Position of n image
            txtbox = slide.shapes.add_textbox(left,top,width=Inches(0.1),height=Inches(0.1)) #add textbox
            p = txtbox.text_frame.paragraphs[0] #access to text, text formatting and paragraph within the textbox
            p.alignment, p.font.size, p.font.bold  = PP_ALIGN.CENTER, Pt(8),True #centers, bold, font size of text
            p.text = letter
        elif n>1 and left<Inches(2) and top==Inches(0.1):
            n,left,top = n+1, left + Inches(2.8),top
            txtbox = slide.shapes.add_textbox(left,top,width=Inches(0.1),height=Inches(0.1)) #add textbox
            p = txtbox.text_frame.paragraphs[0] #access to text, text formatting and paragraph within the textbox
            p.alignment, p.font.size, p.font.bold  = PP_ALIGN.CENTER, Pt(8),True #centers, bold, font size of text
            p.text = letter
        elif n>1 and left>Inches(2) and top<Inches(2):
            n,left,top=n+1,Inches(0.1),top+Inches(1.96)
            f,t = 8,letter
            txtbox = slide.shapes.add_textbox(left,top,width=Inches(0.1),height=Inches(0.1)) #add textbox
            p = txtbox.text_frame.paragraphs[0] #access to text, text formatting and paragraph within the textbox
            p.alignment, p.font.size, p.font.bold  = PP_ALIGN.CENTER, Pt(8),True #centers, bold, font size of text
            p.text = letter
        elif n>1 and left<Inches(2) and top>Inches(1.96):
            n,left,top=n+1,left + Inches(2.8),top
            f,t = 8,letter
            txtbox = slide.shapes.add_textbox(left,top,width=Inches(0.1),height=Inches(0.1)) #add textbox
            p = txtbox.text_frame.paragraphs[0] #access to text, text formatting and paragraph within the textbox
            p.alignment, p.font.size, p.font.bold  = PP_ALIGN.CENTER, Pt(8),True #centers, bold, font size of text
            p.text = letter

    ppt.save('figures.pptx')

    #Export slide to JPG
    Application = win32com.client.Dispatch("PowerPoint.Application") #manipulates powerpoint application directly
    Presentation = Application.Presentations.Open(path+r"\figures.pptx") #open presentation
    Presentation.Slides[0].Export(path + r"\figures.jpg", "JPG") #exports presentation slide
    Application.Quit() #exists presentation
    Presentation =  None
    Application = None